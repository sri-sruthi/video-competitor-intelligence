from __future__ import annotations

import tempfile
from pathlib import Path
import sys
import unittest
import zipfile
from xml.etree import ElementTree as ET

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from pptx import Presentation

from backend.services.pptx_service import PPTXService


def extract_slide_text(slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            parts.append(shape.text)
        if hasattr(shape, "table"):
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text:
                        parts.append(cell.text)
    return "\n".join(parts)


def make_report_payload() -> dict:
    companies = ["HubSpot", "Salesforce"]
    return {
        "companies": companies,
        "report_date": "May 19, 2026",
        "executive_summary": "HubSpot leads on SEO while Salesforce leads on engagement.",
        "company_data": {
            "HubSpot": {
                "subscriber_count": 175000,
                "video_count": 360,
                "published_at": "2015-08-06T14:17:50Z",
                "mean_gap_days": 7.0,
                "consistency_score": 61.5,
                "funnel_label": "Top-Heavy",
                "tofu_pct": 60.0,
                "mofu_pct": 30.0,
                "bofu_pct": 10.0,
                "rolling_avg_engagement": [1.1, 1.3, 1.4, 1.2, 1.5],
                "avg_engagement_rate": 1.42,
                "engagement_trend_slope": 0.02,
                "top_videos": [
                    {
                        "title": "Introducing HubSpot AI",
                        "view_count": 3200,
                        "like_count": 80,
                        "comment_count": 9,
                        "video_id": "abc123",
                        "published_at": "2026-05-01T10:00:00Z",
                        "duration_seconds": 95,
                    },
                    {
                        "title": "How to automate onboarding",
                        "view_count": 2400,
                        "like_count": 44,
                        "comment_count": 5,
                        "video_id": "def456",
                        "published_at": "2026-05-02T10:00:00Z",
                        "duration_seconds": 122,
                    },
                ],
                "short": {"count": 4, "avg_views": 2500.0, "avg_engagement_rate": 1.5},
                "medium": {"count": 1, "avg_views": 1800.0, "avg_engagement_rate": 1.2},
                "long": {"count": 0, "avg_views": 0.0, "avg_engagement_rate": 0.0},
                "best_performing_length": "short",
            },
            "Salesforce": {
                "subscriber_count": 863000,
                "video_count": 1853,
                "published_at": "2006-08-03T21:52:39Z",
                "mean_gap_days": 11.0,
                "consistency_score": 48.2,
                "funnel_label": "Balanced",
                "tofu_pct": 40.0,
                "mofu_pct": 40.0,
                "bofu_pct": 20.0,
                "rolling_avg_engagement": [1.8, 1.9, 2.0, 1.95, 2.1],
                "avg_engagement_rate": 1.95,
                "engagement_trend_slope": -0.01,
                "top_videos": [
                    {
                        "title": "Agentforce platform demo",
                        "view_count": 5100,
                        "like_count": 120,
                        "comment_count": 11,
                        "video_id": "ghi789",
                        "published_at": "2026-05-03T10:00:00Z",
                        "duration_seconds": 88,
                    },
                    {
                        "title": "Enterprise AI workflow guide",
                        "view_count": 3900,
                        "like_count": 62,
                        "comment_count": 8,
                        "video_id": "jkl012",
                        "published_at": "2026-05-04T10:00:00Z",
                        "duration_seconds": 144,
                    },
                ],
                "short": {"count": 3, "avg_views": 4100.0, "avg_engagement_rate": 2.1},
                "medium": {"count": 2, "avg_views": 2500.0, "avg_engagement_rate": 1.6},
                "long": {"count": 0, "avg_views": 0.0, "avg_engagement_rate": 0.0},
                "best_performing_length": "short",
            },
        },
        "rpi_scores": {
            "HubSpot": {"rpi_score": 74.3, "rank": 2},
            "Salesforce": {"rpi_score": 81.8, "rank": 1},
        },
        "topic_clusters": {
            "clusters": [
                {"id": 0, "label": "crm ai / onboarding / automation", "top_terms": [], "size": 3},
                {"id": 1, "label": "agentforce / enterprise ai / demo", "top_terms": [], "size": 2},
            ],
            "company_coverage": {
                "HubSpot": [0],
                "Salesforce": [1],
            },
            "company_theme_labels": {
                "HubSpot": ["crm ai / onboarding / automation"],
                "Salesforce": ["agentforce / enterprise ai / demo"],
            },
            "cluster_coverage": {
                "customer proof / roi / case study": {"HubSpot": 0, "Salesforce": 0},
            },
            "gap_topics": ["customer proof / roi / case study"],
        },
        "gap_analysis": "There is whitespace around ROI proof and customer evidence content.",
        "opportunity_scores": [
            {
                "topic": "customer proof / roi / case study",
                "opportunity_score": 82.4,
                "trend_interest": 64.0,
                "scarcity_score": 100.0,
                "recommendation_type": "Launch new topic cluster",
            }
        ],
        "recommendations": [
            {
                "title": "Create customer proof videos",
                "action": "Ship one case-study-led video per month.",
                "expected_impact": "Better BOFU coverage.",
                "priority": "high",
            },
            {
                "title": "Improve demo coverage",
                "action": "Add shorter demos for high-intent terms.",
                "expected_impact": "More conversion-ready views.",
                "priority": "medium",
            },
        ],
        "action_plan": "Week 1–2: plan. Month 1: publish. Month 2–3: optimize.",
        "seo_scores": {
            "HubSpot": {
                "seo_score": 78.9,
                "breakdown": {
                    "title_length_score": 95.6,
                    "description_depth": 100.0,
                    "has_timestamps_pct": 0.0,
                    "tags_count_avg": 11.8,
                    "keyword_in_title_score": 100.0,
                },
            },
            "Salesforce": {
                "seo_score": 66.7,
                "breakdown": {
                    "title_length_score": 80.0,
                    "description_depth": 100.0,
                    "has_timestamps_pct": 0.0,
                    "tags_count_avg": 6.0,
                    "keyword_in_title_score": 80.0,
                },
            },
        },
        "slide_interpretations": {
            "slide_03": "Custom LLM interpretation for the channel overview slide.",
            "slide_13": "Custom LLM interpretation for the high-priority moves slide.",
        },
    }


class PPTXServiceTests(unittest.TestCase):
    def setUp(self) -> None:
        self.service = PPTXService()

    def test_generate_report_creates_16_slide_pptx(self) -> None:
        payload = make_report_payload()

        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = str(Path(tmpdir) / "phase3-report.pptx")
            result = self.service.generate_report(payload, output_path)

            self.assertEqual(result, output_path)
            self.assertTrue(Path(output_path).exists())

            prs = Presentation(output_path)
            self.assertEqual(len(prs.slides), 16)

    def test_generated_deck_contains_expected_headings_and_company_text(self) -> None:
        payload = make_report_payload()

        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = str(Path(tmpdir) / "phase3-report.pptx")
            self.service.generate_report(payload, output_path)
            prs = Presentation(output_path)

            cover_text = extract_slide_text(prs.slides[0])
            summary_text = extract_slide_text(prs.slides[1])
            overview_text = extract_slide_text(prs.slides[2])
            recommendations_text = extract_slide_text(prs.slides[12])
            scorecard_text = extract_slide_text(prs.slides[14])

            self.assertIn("Client Growth", cover_text)
            self.assertIn("HubSpot vs Salesforce", cover_text)
            self.assertIn("Executive Summary", summary_text)
            self.assertIn("Salesforce", summary_text)
            self.assertIn("Company Growth Scorecard", scorecard_text)
            self.assertIn("HubSpot", scorecard_text)
            self.assertIn("Salesforce", scorecard_text)
            self.assertIn("What You'll Learn", cover_text)
            self.assertIn("Current gap to the next-ranked peer", summary_text)
            self.assertIn("Custom LLM interpretation for the channel overview slide.", overview_text)
            self.assertIn("Custom LLM interpretation for the high-priority moves slide.", recommendations_text)
            self.assertIn("Grade bands", scorecard_text)

    def test_generated_deck_keeps_shapes_within_slide_bounds_and_formats_small_gaps_humanly(self) -> None:
        payload = make_report_payload()
        payload["company_data"]["HubSpot"]["mean_gap_days"] = 0.02
        payload["recommendations"][0]["data_rationale"] = "BOFU coverage is underweight."

        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = str(Path(tmpdir) / "phase3-report.pptx")
            self.service.generate_report(payload, output_path)

            prs = Presentation(output_path)
            cadence_text = extract_slide_text(prs.slides[4])
            funnel_text = extract_slide_text(prs.slides[5])
            topics_text = extract_slide_text(prs.slides[8])
            gap_text = extract_slide_text(prs.slides[9])
            recommendation_text = extract_slide_text(prs.slides[12])
            self.assertIn("<1 day", cadence_text)
            self.assertIn("Burst / same-day", cadence_text)
            self.assertIn("confidence", funnel_text.lower())
            self.assertIn("Representative themes", topics_text)
            self.assertIn("Where To Test Next", gap_text)
            self.assertIn("Why this matters", recommendation_text)
            self.assertIn("Business impact", recommendation_text)

            ns = {
                "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
                "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            }
            with zipfile.ZipFile(output_path) as archive:
                root = ET.fromstring(archive.read("ppt/presentation.xml"))
                slide_size = root.find("p:sldSz", ns)
                self.assertIsNotNone(slide_size)
                slide_width = int(slide_size.attrib["cx"])
                slide_height = int(slide_size.attrib["cy"])

                for slide_number in range(1, 17):
                    xml = ET.fromstring(archive.read(f"ppt/slides/slide{slide_number}.xml"))
                    for xfrm in xml.findall(".//a:xfrm", ns):
                        off = xfrm.find("a:off", ns)
                        ext = xfrm.find("a:ext", ns)
                        if off is None or ext is None:
                            continue
                        x = int(off.attrib.get("x", 0))
                        y = int(off.attrib.get("y", 0))
                        width = int(ext.attrib.get("cx", 0))
                        height = int(ext.attrib.get("cy", 0))
                        self.assertGreaterEqual(x, 0, f"slide {slide_number} x below 0")
                        self.assertGreaterEqual(y, 0, f"slide {slide_number} y below 0")
                        self.assertLessEqual(
                            x + width,
                            slide_width,
                            f"slide {slide_number} shape exceeds right edge",
                        )
                        self.assertLessEqual(
                            y + height,
                            slide_height,
                            f"slide {slide_number} shape exceeds bottom edge",
                        )

            with zipfile.ZipFile(output_path) as archive:
                rels_root = ET.fromstring(
                    archive.read("ppt/slides/_rels/slide8.xml.rels")
                )
                targets = [rel.attrib.get("Target", "") for rel in rels_root]
                self.assertTrue(
                    any("youtube.com/watch?v=" in target for target in targets),
                    "slide 8 should contain hyperlink relationships for videos",
                )


if __name__ == "__main__":
    unittest.main()
