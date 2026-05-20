from __future__ import annotations

import json
import tempfile
from pathlib import Path
import sys
import unittest
from unittest.mock import AsyncMock, patch

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from pptx import Presentation

from backend.services.analytics_service import AnalyticsService
from backend.services.pptx_service import PPTXService
from tests.phase2_test_support import load_ai_module, load_seo_module


AI_MODULE, FakeClient = load_ai_module()
SEO_MODULE, FakeTrendReq = load_seo_module()
AIService = AI_MODULE.AIService
SEOService = SEO_MODULE.SEOService


def make_videos(prefix: str) -> list[dict]:
    return [
        {
            "title": f"Introducing {prefix} AI workflows",
            "description": "00:00 Intro\n" + ("x" * 320),
            "tags": ["ai", "crm", "workflows", "sales"],
            "published_at": "2026-05-01T00:00:00Z",
            "view_count": 1200,
            "like_count": 30,
            "comment_count": 4,
            "duration_seconds": 210,
        },
        {
            "title": f"How to launch {prefix} onboarding faster",
            "description": "00:00 Setup\n" + ("x" * 280),
            "tags": ["onboarding", "tutorial", "saas"],
            "published_at": "2026-05-05T00:00:00Z",
            "view_count": 850,
            "like_count": 22,
            "comment_count": 6,
            "duration_seconds": 430,
        },
        {
            "title": f"The future of {prefix} customer education",
            "description": "Perspective video " + ("y" * 210),
            "tags": ["thought leadership", "education"],
            "published_at": "2026-05-10T00:00:00Z",
            "view_count": 660,
            "like_count": 11,
            "comment_count": 2,
            "duration_seconds": 620,
        },
        {
            "title": f"5 ways {prefix} demos convert better",
            "description": "00:00 Tips\n" + ("z" * 305),
            "tags": ["demos", "conversion", "listicle", "video"],
            "published_at": "2026-05-17T00:00:00Z",
            "view_count": 980,
            "like_count": 25,
            "comment_count": 5,
            "duration_seconds": 260,
        },
        {
            "title": f"{prefix} integration walkthrough for ops teams",
            "description": "00:00 Walkthrough\n" + ("w" * 330),
            "tags": ["integration", "ops", "tutorial", "implementation"],
            "published_at": "2026-05-24T00:00:00Z",
            "view_count": 720,
            "like_count": 19,
            "comment_count": 3,
            "duration_seconds": 540,
        },
    ]


class Phase3IntegrationTests(unittest.TestCase):
    def setUp(self) -> None:
        FakeClient.reset()
        FakeTrendReq.reset()
        self.analytics = AnalyticsService()
        self.seo = SEOService()
        self.ai = AIService("test-key")
        self.pptx = PPTXService()

    def _build_company_record(
        self,
        company_name: str,
        videos: list[dict],
        *,
        subscriber_count: int,
        video_count: int,
        published_at: str,
    ) -> tuple[dict, dict]:
        engagement = self.analytics.compute_engagement_metrics(videos)
        consistency = self.analytics.compute_upload_consistency(videos)
        funnel = self.analytics.classify_content_funnel(videos)
        length = self.analytics.analyse_video_length_strategy(videos)
        format_mix = self.analytics.analyse_format_mix(videos)
        recent = self.analytics.compute_recent_view_velocity(videos)
        titles = self.analytics.analyse_title_patterns(videos)
        seo_score = self.seo.score_video_seo(videos)

        company_row = {
            "subscriber_count": subscriber_count,
            "video_count": video_count,
            "published_at": published_at,
            **engagement,
            **consistency,
            **funnel,
            **length,
            **format_mix,
            **recent,
            "dominant_strategy_label": titles["dominant_strategy"]["label"],
            "dominant_strategy_description": titles["dominant_strategy"]["description"],
            "tutorial_count": titles["tutorial_count"],
            "listicle_count": titles["listicle_count"],
            "thought_leadership_count": titles["thought_leadership_count"],
            "product_led_count": titles["product_led_count"],
            "top_videos": sorted(videos, key=lambda v: v.get("view_count", 0), reverse=True)[:2],
        }

        rpi_seed = {
            "avg_engagement_rate": engagement["avg_engagement_rate"],
            "views_per_subscriber": round(
                engagement["views_per_video"] / max(subscriber_count, 1), 6
            ),
            "consistency_score": consistency["consistency_score"],
            "topic_diversity_score": max(
                titles["tutorial_count"]
                + titles["listicle_count"]
                + titles["thought_leadership_count"]
                + titles["product_led_count"],
                1,
            ) * 25.0,
            "seo_score": seo_score["seo_score"],
        }
        return company_row, {"seo": seo_score, "rpi_seed": rpi_seed}

    def test_phase1_phase2_phase3_pipeline_builds_a_real_pptx_report(self) -> None:
        companies = ["HubSpot", "Salesforce"]
        videos_by_company = {
            "HubSpot": make_videos("HubSpot"),
            "Salesforce": make_videos("Salesforce"),
        }
        channel_meta = {
            "HubSpot": {
                "subscriber_count": 175000,
                "video_count": 360,
                "published_at": "2015-08-06T14:17:50Z",
            },
            "Salesforce": {
                "subscriber_count": 863000,
                "video_count": 1853,
                "published_at": "2006-08-03T21:52:39Z",
            },
        }

        company_data: dict[str, dict] = {}
        seo_scores: dict[str, dict] = {}
        rpi_input: dict[str, dict] = {}

        for company in companies:
            company_row, aux = self._build_company_record(
                company,
                videos_by_company[company],
                **channel_meta[company],
            )
            company_data[company] = company_row
            seo_scores[company] = aux["seo"]
            rpi_input[company] = aux["rpi_seed"]

        topic_clusters = self.analytics.cluster_content_topics(videos_by_company)
        trends = {
            topic: {
                "avg_interest": 65.0 - index * 7,
                "trend_direction": "rising" if index == 0 else "flat",
                "peak_month": "Apr",
                "fallback_used": False,
                "source": "test",
                "retries_used": 0,
                "error_reason": "",
            }
            for index, topic in enumerate(topic_clusters["gap_topics"])
        }
        opportunity_scores = self.seo.compute_opportunity_scores(
            topic_clusters["gap_topics"],
            trends,
            topic_clusters["cluster_coverage"],
        )
        rpi_scores = self.analytics.compute_rpi(rpi_input)

        FakeClient.queued_results = [
            "Executive summary for the deck",
            "Gap analysis for the deck",
            json.dumps(
                [
                    {
                        "title": f"Recommendation {i}",
                        "action": "Publish deeper conversion content",
                        "data_rationale": "The funnel remains awareness-heavy.",
                        "expected_impact": "Stronger mid-funnel engagement.",
                        "priority": "high" if i == 1 else "medium",
                    }
                    for i in range(1, 6)
                ]
            ),
            "90-day action plan for the deck",
        ]

        with patch.object(AI_MODULE.asyncio, "sleep", new=AsyncMock()):
            executive_summary = self.ai.generate_executive_summary(
                {
                    company: {
                        **company_data[company],
                        "rpi_score": rpi_scores[company]["rpi_score"],
                        "rank": rpi_scores[company]["rank"],
                    }
                    for company in companies
                }
            )
            gap_analysis = self.ai.generate_gap_analysis(topic_clusters["gap_topics"], trends)
            recommendations = self.ai.generate_recommendations(
                "HubSpot",
                {
                    **company_data["HubSpot"],
                    "seo_score": seo_scores["HubSpot"]["seo_score"],
                    "rpi_score": rpi_scores["HubSpot"]["rpi_score"],
                    "rank": rpi_scores["HubSpot"]["rank"],
                },
                topic_clusters["gap_topics"],
            )
            action_plan = self.ai.generate_action_plan(
                recommendations,
                current_cadence=company_data["HubSpot"]["mean_gap_days"],
            )

        payload = {
            "companies": companies,
            "report_date": "May 19, 2026",
            "executive_summary": executive_summary,
            "company_data": company_data,
            "rpi_scores": rpi_scores,
            "topic_clusters": topic_clusters,
            "gap_analysis": gap_analysis,
            "opportunity_scores": opportunity_scores,
            "recommendations": recommendations,
            "action_plan": action_plan,
            "seo_scores": seo_scores,
        }

        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = str(Path(tmpdir) / "integration-report.pptx")
            self.pptx.generate_report(payload, output_path)

            prs = Presentation(output_path)
            self.assertEqual(len(prs.slides), 16)

            all_text = "\n".join(
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text
            )

            self.assertIn("Executive Summary", all_text)
            self.assertIn("Executive brief", all_text)
            self.assertIn("HubSpot", all_text)
            self.assertIn("Salesforce", all_text)
            self.assertIn("Recommendation 1", all_text)
            self.assertIn("Video Marketing Health Score", all_text)


if __name__ == "__main__":
    unittest.main()
