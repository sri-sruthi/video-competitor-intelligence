"""
pptx_service.py
Generates a 16-slide professional PowerPoint report for the
Video Competitor Intelligence tool.

Dependencies: python-pptx, matplotlib
Compatible with all Phase 1 (youtube_service, analytics_service) and
Phase 2 (ai_service, seo_service) return structures.
"""

from __future__ import annotations

import io
import math
import re
from datetime import datetime
from typing import Any

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE


# ---------------------------------------------------------------------------
# Design system constants
# ---------------------------------------------------------------------------

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Colours (RGB tuples — used for python-pptx)
C_NAVY = RGBColor(20, 33, 51)
C_WHITE = RGBColor(255, 255, 255)
C_BLUE = RGBColor(37, 99, 235)
C_TEAL = RGBColor(20, 184, 166)
C_DARK = RGBColor(30, 41, 59)
C_GREY = RGBColor(100, 116, 139)
C_LIGHT_BG = RGBColor(247, 244, 238)
C_SOFT_PANEL = RGBColor(237, 242, 247)
C_AMBER = RGBColor(245, 158, 11)
C_GREEN = RGBColor(34, 197, 94)
C_RED = RGBColor(239, 68, 68)
C_MUTED_TEAL = RGBColor(11, 122, 119)
C_BORDER = RGBColor(220, 227, 234)
C_PAPER = RGBColor(252, 250, 246)

# Matplotlib hex equivalents (no # — consistent across charts)
MPL_NAVY   = "#142133"
MPL_BLUE   = "#2563EB"
MPL_TEAL   = "#14B8A6"
MPL_AMBER  = "#F59E0B"
MPL_GREEN  = "#22C55E"
MPL_RED    = "#EF4444"
MPL_PURPLE = "#8B5CF6"
MPL_PINK   = "#EC4899"

# Per-company colour cycle (used consistently on all charts)
COMPANY_COLOURS = [MPL_BLUE, MPL_TEAL, MPL_AMBER, MPL_GREEN, MPL_RED]

FONT_HEADING = "Aptos Display"
FONT_BODY    = "Aptos"
FONT_COVER   = "Georgia"

FOOTER_TEXT = "Client Growth Intelligence Report"


# ---------------------------------------------------------------------------
# PPTXService
# ---------------------------------------------------------------------------

class PPTXService:
    """Generates a 16-slide Video Competitor Intelligence PowerPoint report."""

    # ------------------------------------------------------------------
    # Public entry point
    # ------------------------------------------------------------------

    def generate_report(self, data: dict, output_path: str) -> str:
        """
        Build the full 16-slide PPTX report and write it to output_path.

        Args:
            data: Dict containing all report data. Expected keys:
                  companies (list[str]),
                  report_date (str),
                  executive_summary (str),
                  company_data (dict[str, dict]),   ← merged analytics per company
                  rpi_scores (dict),                ← from compute_rpi()
                  topic_clusters (dict),            ← from cluster_content_topics()
                  gap_analysis (str),
                  opportunity_scores (list[dict]),  ← from compute_opportunity_scores()
                  recommendations (list[dict]),
                  action_plan (str),
                  seo_scores (dict[str, dict]),     ← from score_video_seo()
            output_path: Filesystem path for the .pptx file.

        Returns:
            The output_path string.
        """
        prs = Presentation()
        prs.slide_width  = SLIDE_W
        prs.slide_height = SLIDE_H

        companies    = data.get("companies", [])
        report_date  = data.get("report_date", datetime.today().strftime("%B %d, %Y"))
        company_data = data.get("company_data", {})
        rpi_scores   = data.get("rpi_scores", {})
        topic_clusters = data.get("topic_clusters", {})
        gap_analysis   = data.get("gap_analysis", "")
        opp_scores     = data.get("opportunity_scores", [])
        recommendations = data.get("recommendations", [])
        action_plan    = data.get("action_plan", "")
        seo_scores     = data.get("seo_scores", {})
        exec_summary   = data.get("executive_summary", "")
        self._active_interpretations = data.get("slide_interpretations", {}) or {}

        # Build company colour map (stable across slides)
        company_colours = {
            c: COMPANY_COLOURS[i % len(COMPANY_COLOURS)]
            for i, c in enumerate(companies)
        }

        blank = prs.slide_layouts[6]  # completely blank layout

        self._slide_01_cover(prs.slides.add_slide(blank), companies, report_date)
        self._slide_02_exec_summary(prs.slides.add_slide(blank), exec_summary, rpi_scores, company_data, seo_scores)
        target_company = data.get("target_company", companies[0] if companies else "Your Company")

        self._slide_03_channel_overview(prs.slides.add_slide(blank), companies, company_data, company_colours, target_company)
        self._slide_04_rpi_rankings(prs.slides.add_slide(blank), companies, company_data, rpi_scores, company_colours, target_company)
        self._slide_05_upload_cadence(prs.slides.add_slide(blank), companies, company_data, company_colours)
        self._slide_06_funnel(prs.slides.add_slide(blank), companies, company_data, company_colours)
        self._slide_07_engagement(prs.slides.add_slide(blank), companies, company_data, company_colours)
        self._slide_08_top_videos(prs.slides.add_slide(blank), companies, company_data, company_colours)
        self._slide_09_topics(prs.slides.add_slide(blank), companies, topic_clusters, company_colours)
        self._slide_12_gap_analysis(prs.slides.add_slide(blank), gap_analysis, opp_scores)
        self._slide_10_length_strategy(prs.slides.add_slide(blank), companies, company_data, company_colours)
        self._slide_11_seo(prs.slides.add_slide(blank), companies, seo_scores, company_data, company_colours)
        self._slide_13_priority_moves(prs.slides.add_slide(blank), recommendations, "high")
        self._slide_13_priority_moves(prs.slides.add_slide(blank), recommendations, "medium")
        self._slide_14_scorecard(prs.slides.add_slide(blank), companies, company_data, rpi_scores, seo_scores, report_date, target_company)
        self._slide_15_action_plan(prs.slides.add_slide(blank), action_plan, recommendations, target_company, company_data.get(target_company, {}), opp_scores)

        prs.save(output_path)
        return output_path

    # ==================================================================
    # Slide builders
    # ==================================================================

    # ------------------------------------------------------------------
    # SLIDE 1 — Cover (light theme)
    # ------------------------------------------------------------------

    def _slide_01_cover(self, slide, companies: list[str], report_date: str):
        self._set_bg(slide, C_PAPER)

        # Left accent rail and scope panel
        self._rect(slide, 0, 0, Inches(0.16), SLIDE_H, C_MUTED_TEAL)
        self._rect(slide, Inches(9.35), Inches(0.4), Inches(3.45), Inches(2.05), C_SOFT_PANEL)
        self._text_box(
            slide, "What You'll Learn",
            Inches(9.6), Inches(0.7), Inches(2.8), Inches(0.3),
            font_size=12, bold=True, color=C_MUTED_TEAL,
        )
        self._text_box(
            slide,
            f"• {len(companies)} peer companies\n"
            "• Audience, cadence, and discovery signals\n"
            "• Clear growth actions for the client team",
            Inches(9.6), Inches(1.05), Inches(2.8), Inches(1.05),
            font_size=11, color=C_DARK, wrap=True,
        )

        self._text_box(
            slide, "How your video strategy compares and where to grow",
            Inches(0.6), Inches(0.8), Inches(5.2), Inches(0.35),
            font_size=12, bold=True, color=C_MUTED_TEAL, align=PP_ALIGN.LEFT,
            font_face=FONT_BODY,
        )

        self._text_box(
            slide,
            "Client Growth\nIntelligence",
            Inches(0.6), Inches(1.4), Inches(9.0), Inches(2.4),
            font_size=48, bold=True, color=C_NAVY, align=PP_ALIGN.LEFT,
            font_face=FONT_COVER,
        )

        self._text_box(
            slide,
            "Video strategy benchmark for executive decision-making",
            Inches(0.62), Inches(3.42), Inches(6.2), Inches(0.34),
            font_size=15, color=C_GREY, align=PP_ALIGN.LEFT,
        )

        subtitle = " vs ".join(companies) if companies else "Competitive Analysis"
        self._text_box(
            slide, subtitle,
            Inches(0.6), Inches(4.15), Inches(10.0), Inches(0.7),
            font_size=21, bold=True, color=C_MUTED_TEAL, align=PP_ALIGN.LEFT,
        )

        self._text_box(
            slide, report_date,
            Inches(0.6), Inches(4.85), Inches(6.0), Inches(0.45),
            font_size=14, color=C_GREY, align=PP_ALIGN.LEFT,
        )

        self._rect(slide, Inches(0.6), Inches(3.75), Inches(6.15), Inches(0.04), C_MUTED_TEAL)
        self._footer(slide, dark=False)

    # ------------------------------------------------------------------
    # SLIDE 2 — Executive Summary (white)
    # ------------------------------------------------------------------

    def _slide_02_exec_summary(self, slide, exec_summary: str, rpi_scores: dict, company_data: dict, seo_scores: dict):
        self._set_bg(slide, C_WHITE)
        self._slide_heading(slide, "Executive Summary")

        summary_source = self._slide_insight("slide_02", exec_summary or "No summary generated.")
        intro, model_bullets, closing = self._executive_summary_parts(summary_source)
        bullets = model_bullets or self._executive_bullets(company_data, rpi_scores, seo_scores)
        self._rect(slide, Inches(0.45), Inches(1.28), Inches(8.15), Inches(4.95), C_LIGHT_BG)
        self._text_box(
            slide, "Executive brief",
            Inches(0.7), Inches(1.55), Inches(2.0), Inches(0.35),
            font_size=12, bold=True, color=C_MUTED_TEAL,
        )
        self._text_box(
            slide, intro,
            Inches(0.7), Inches(1.92), Inches(7.2), Inches(0.85),
            font_size=15, bold=False, color=C_DARK, wrap=True, fit=True, min_font_size=12,
        )
        bullet_y = 2.95
        for bullet in bullets[:4]:
            self._text_box(
                slide, f"• {bullet}",
                Inches(0.88), Inches(bullet_y), Inches(7.0), Inches(0.4),
                font_size=12, color=C_DARK, wrap=True, fit=True, min_font_size=10,
            )
            bullet_y += 0.52

        self._text_box(
            slide, "What this means",
            Inches(0.72), Inches(5.15), Inches(2.0), Inches(0.25),
            font_size=11, bold=True, color=C_MUTED_TEAL,
        )
        self._text_box(
            slide, closing,
            Inches(0.72), Inches(5.43), Inches(7.55), Inches(0.6),
            font_size=11, color=C_DARK, wrap=True, fit=True, min_font_size=9,
        )

        if rpi_scores:
            top_company = min(rpi_scores, key=lambda c: rpi_scores[c].get("rank", 99))
            top_rpi     = rpi_scores[top_company].get("rpi_score", 0)
            spread = 0.0
            ranked_vals = sorted(
                (payload.get("rpi_score", 0.0) for payload in rpi_scores.values()),
                reverse=True,
            )
            if len(ranked_vals) >= 2:
                spread = ranked_vals[0] - ranked_vals[1]

            self._rect(slide, Inches(9.0), Inches(1.3), Inches(3.8), Inches(2.6), C_NAVY)
            self._text_box(
                slide, "#1 RANKED",
                Inches(9.0), Inches(1.4), Inches(3.8), Inches(0.5),
                font_size=12, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER,
            )
            self._text_box(
                slide, top_company,
                Inches(9.0), Inches(1.95), Inches(3.8), Inches(0.9),
                font_size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
            )
            self._text_box(
                slide, f"RPI Score: {top_rpi:.1f}",
                Inches(9.0), Inches(2.9), Inches(3.8), Inches(0.5),
                font_size=16, color=C_AMBER, align=PP_ALIGN.CENTER,
            )
            self._rect(slide, Inches(9.0), Inches(4.1), Inches(3.8), Inches(1.7), C_SOFT_PANEL)
            self._text_box(
                slide,
                (
                    f"Current gap to the next-ranked peer: {spread:.1f} points.\n"
                    "Use the following slides to see whether that lead comes from audience pull, operating rhythm, discovery setup, or better buyer-stage coverage."
                ),
            Inches(9.2), Inches(4.3), Inches(3.35), Inches(1.42),
            font_size=10, color=C_DARK, wrap=True,
        )

        self._footer(slide)

    # ------------------------------------------------------------------
    # SLIDE 3 — Channel Overview (white)
    # ------------------------------------------------------------------

    def _slide_03_channel_overview(self, slide, companies, company_data, company_colours, target_company: str):
        self._set_bg(slide, C_WHITE)
        self._slide_heading(slide, "Channel Overview")

        subs = [company_data.get(c, {}).get("subscriber_count", 0) for c in companies]
        colours = [company_colours[c] for c in companies]
        chart_img = self._bar_chart_h(
            labels=companies,
            values=subs,
            colours=colours,
            title="Subscribers by Company",
            xlabel="Subscribers",
        )
        slide.shapes.add_picture(chart_img, Inches(0.4), Inches(1.2), Inches(6.9), Inches(3.0))

        if companies:
            largest = max(companies, key=lambda c: company_data.get(c, {}).get("subscriber_count", 0))
            largest_subs = company_data.get(largest, {}).get("subscriber_count", 0)
            insight_fallback = self._channel_overview_takeaway(companies, company_data, target_company)
            insight = self._slide_insight(
                "slide_03",
                insight_fallback,
            )
            self._rect(slide, Inches(7.55), Inches(1.35), Inches(5.25), Inches(2.8), C_SOFT_PANEL)
            self._text_box(
                slide, "Why this matters",
                Inches(7.8), Inches(1.58), Inches(2.3), Inches(0.32),
                font_size=12, bold=True, color=C_MUTED_TEAL,
            )
            self._text_box(
                slide, insight,
                Inches(7.8), Inches(1.95), Inches(4.75), Inches(1.95),
                font_size=11, color=C_DARK, wrap=True, fit=True, min_font_size=9,
            )

        headers = ["Company", "Subscribers", "Total Videos", "Channel Age (yrs)"]
        rows = []
        for c in companies:
            cd = company_data.get(c, {})
            subs_val  = f"{cd.get('subscriber_count', 0):,}"
            vids_val  = f"{cd.get('video_count', 0):,}"
            pub       = cd.get("published_at", "")
            age       = self._channel_age_years(pub)
            rows.append([c, subs_val, vids_val, age])

        self._table(slide, headers, rows, Inches(0.4), Inches(4.45), Inches(12.4), Inches(2.45))
        self._footer(slide)

    # ------------------------------------------------------------------
    # SLIDE 4 — RPI Rankings (white)
    # ------------------------------------------------------------------

    def _slide_04_rpi_rankings(self, slide, companies, company_data, rpi_scores, company_colours, target_company: str):
        self._set_bg(slide, C_WHITE)
        self._slide_heading(slide, "Video Marketing Health Score")
        self._text_box(
            slide,
            "Composite benchmark across engagement, SEO, cadence, content breadth, and audience response",
            Inches(0.45), Inches(0.78), Inches(8.0), Inches(0.22),
            font_size=10, color=C_GREY,
        )

        sorted_companies = sorted(rpi_scores, key=lambda c: rpi_scores[c].get("rank", 99))
        values  = [rpi_scores[c].get("rpi_score", 0) for c in sorted_companies]
        colours = [company_colours.get(c, MPL_BLUE) for c in sorted_companies]

        chart_img = self._bar_chart_h(
            labels=sorted_companies,
            values=values,
            colours=colours,
            title="RPI Score (0–100)",
            xlabel="Health score",
            xlim=(0, 100),
        )
        slide.shapes.add_picture(chart_img, Inches(0.4), Inches(1.2), Inches(8.5), Inches(4.0))

        ranked = sorted(
            ((company, payload.get("rpi_score", 0.0)) for company, payload in rpi_scores.items()),
            key=lambda item: item[1],
            reverse=True,
        )
        leader = ranked[0][0] if ranked else "—"
        benchmark_company, benchmark_reason = self._benchmark_peer_explanation(
            companies, company_data, rpi_scores, target_company
        )
        explanation = (
            "This score is a board-friendly summary, not a raw business KPI.\n\n"
            "It blends five signals into one benchmark so a non-technical reader can quickly see who appears strongest overall:\n"
            "• audience response\n"
            "• reach efficiency\n"
            "• publishing discipline\n"
            "• content breadth\n"
            "• discoverability"
        )
        self._rect(slide, Inches(9.0), Inches(1.45), Inches(3.9), Inches(2.95), C_LIGHT_BG)
        self._text_box(
            slide, explanation,
            Inches(9.2), Inches(1.65), Inches(3.45), Inches(2.55),
            font_size=12, color=C_DARK, wrap=True,
        )
        self._rect(slide, Inches(9.0), Inches(4.6), Inches(3.9), Inches(1.1), C_SOFT_PANEL)
        self._text_box(
            slide,
            self._slide_insight(
                "slide_04",
                benchmark_reason or f"{leader} leads this peer set. Use the following slides to see which operating habits are most transferable to {target_company}.",
            ),
            Inches(9.2), Inches(4.78), Inches(3.45), Inches(0.86),
            font_size=10, color=C_DARK, wrap=True, fit=True, min_font_size=8,
        )
        self._footer(slide)

    # ------------------------------------------------------------------
    # SLIDE 5 — Upload Cadence & Consistency (white)
    # ------------------------------------------------------------------

    def _slide_05_upload_cadence(self, slide, companies, company_data, company_colours):
        self._set_bg(slide, C_WHITE)
        self._slide_heading(slide, "Upload Cadence & Consistency")

        gaps   = [company_data.get(c, {}).get("mean_gap_days", 0) for c in companies]
        scores = [company_data.get(c, {}).get("consistency_score", 0) for c in companies]
        colours = [company_colours[c] for c in companies]

        # Two side-by-side charts
        gap_img = self._bar_chart_v(
            labels=companies,
            values=gaps,
            colours=colours,
            title="Avg Days Between Uploads (lower = more frequent)",
            ylabel="Days",
            figsize=(5.5, 3.0),
        )
        cons_img = self._bar_chart_v(
            labels=companies,
            values=scores,
            colours=colours,
            title="Consistency Score (0–100)",
            ylabel="Score",
            ylim=(0, 100),
            figsize=(5.5, 3.0),
        )
        slide.shapes.add_picture(gap_img,  Inches(0.3), Inches(1.2), Inches(6.2), Inches(3.0))
        slide.shapes.add_picture(cons_img, Inches(6.8), Inches(1.2), Inches(6.2), Inches(3.0))

        if companies:
            steadiest = max(companies, key=lambda c: company_data.get(c, {}).get("consistency_score", 0))
            burstiest = min(companies, key=lambda c: company_data.get(c, {}).get("consistency_score", 0))
            insight = self._slide_insight(
                "slide_05",
                (
                f"{steadiest} is the steadiest recent publisher in this sample, while {burstiest} shows the most uneven rhythm. "
                "Read any cadence conclusion through the confidence level in the table below, especially when the recent sample is thin."
                ),
            )
            self._text_box(
                slide, insight,
                Inches(0.45), Inches(4.08), Inches(12.2), Inches(0.38),
                font_size=10, color=C_GREY, wrap=True,
            )

        headers = ["Company", "Avg Gap (days)", "Consistency Score", "Cadence", "Confidence"]
        rows = []
        for c in companies:
            cd       = company_data.get(c, {})
            gap_d    = cd.get("mean_gap_days", 0)
            cons_s   = cd.get("consistency_score", 0)
            cadence  = self._cadence_label(gap_d)
            confidence = cd.get("cadence_confidence", "LOW")
            rows.append([c, self._format_gap_days(gap_d), f"{cons_s:.1f}", cadence, confidence.title()])

        self._table(slide, headers, rows, Inches(0.3), Inches(4.45), Inches(12.6), Inches(2.7))
        self._footer(slide)

    # ------------------------------------------------------------------
    # SLIDE 6 — Content Funnel Strategy (white)
    # ------------------------------------------------------------------

    def _slide_06_funnel(self, slide, companies, company_data, company_colours):
        self._set_bg(slide, C_WHITE)
        self._slide_heading(slide, "Audience Journey Coverage")

        tofu_vals = [company_data.get(c, {}).get("tofu_pct", 0) for c in companies]
        mofu_vals = [company_data.get(c, {}).get("mofu_pct", 0) for c in companies]
        bofu_vals = [company_data.get(c, {}).get("bofu_pct", 0) for c in companies]
        unclassified_vals = [company_data.get(c, {}).get("unclassified_pct", 0) for c in companies]

        chart_img = self._stacked_bar_chart(
            companies=companies,
            series={
                "Awareness (TOFU)": tofu_vals,
                "Consideration (MOFU)": mofu_vals,
                "Proof/Conversion (BOFU)": bofu_vals,
                "Unclear intent": unclassified_vals,
            },
            colours=[MPL_TEAL, MPL_BLUE, MPL_AMBER, "#94A3B8"],
            title="Content Funnel Distribution (%)",
        )
        slide.shapes.add_picture(chart_img, Inches(0.4), Inches(1.2), Inches(8.6), Inches(3.75))

        self._rect(slide, Inches(9.02), Inches(1.16), Inches(3.88), Inches(5.62), C_SOFT_PANEL)
        self._text_box(
            slide, "Why this matters",
            Inches(9.24), Inches(1.36), Inches(3.05), Inches(0.32),
            font_size=12, bold=True, color=C_MUTED_TEAL,
        )
        self._text_box(
            slide,
            self._slide_insight(
                "slide_06",
                self._funnel_takeaway(companies, company_data),
            ),
            Inches(9.24), Inches(1.72), Inches(3.28), Inches(0.9),
            font_size=9, color=C_DARK, wrap=True, fit=True, min_font_size=8,
        )
        y_start = 2.72
        for c in companies:
            fl = company_data.get(c, {}).get("funnel_label", "—")
            if fl == "Balanced":
                colour = C_TEAL
            elif fl == "Awareness-focused":
                colour = C_BLUE
            elif fl == "Conversion-focused":
                colour = C_AMBER
            else:
                colour = C_GREY
            unclassified = company_data.get(c, {}).get("unclassified_count", 0)
            confidence = company_data.get(c, {}).get("funnel_confidence", "LOW")
            self._text_box(
                slide, f"{c[:18]}",
                Inches(9.24), Inches(y_start), Inches(3.25), Inches(0.22),
                font_size=11, bold=True, color=C_DARK,
            )
            self._text_box(
                slide, f"{fl} • {unclassified} unclear • {confidence.title()} confidence",
                Inches(9.24), Inches(y_start + 0.2), Inches(3.25), Inches(0.2),
                font_size=9, color=colour, fit=True, min_font_size=8,
            )
            y_start += 0.62

        self._footer(slide)

    # ------------------------------------------------------------------
    # SLIDE 7 — Engagement Rate Analysis (white)
    # ------------------------------------------------------------------

    def _slide_07_engagement(self, slide, companies, company_data, company_colours):
        self._set_bg(slide, C_WHITE)
        self._slide_heading(slide, "Engagement Rate Analysis")

        # Rolling avg engagement line chart
        fig, ax = plt.subplots(figsize=(8, 3.5))
        ax.set_facecolor("white")
        fig.patch.set_facecolor("white")
        for spine in ax.spines.values():
            spine.set_visible(False)
        ax.yaxis.grid(True, color="#E2E8F0", linewidth=0.8)
        ax.set_axisbelow(True)
        ax.tick_params(colors="#64748B")

        for c in companies:
            rolling = company_data.get(c, {}).get("rolling_avg_engagement", [])
            if rolling:
                # Show last 20 points
                data_slice = rolling[-20:]
                ax.plot(
                    range(len(data_slice)),
                    data_slice,
                    label=c,
                    color=company_colours[c],
                    linewidth=2.5,
                    marker="o",
                    markersize=4,
                )

        ax.set_title("Rolling Avg Engagement (last 20 videos)", fontsize=11, color="#1E293B", pad=8)
        ax.set_xlabel("Video Index (recent →)", fontsize=9, color="#64748B")
        ax.set_ylabel("Engagement Rate (%)", fontsize=9, color="#64748B")
        ax.legend(fontsize=9, frameon=False)
        plt.tight_layout()
        chart_img = self._fig_to_image(fig)

        slide.shapes.add_picture(chart_img, Inches(0.4), Inches(1.2), Inches(8.5), Inches(3.8))

        # Table (right)
        headers = ["Company", "Avg ER (%)", "Trend"]
        rows = []
        for c in companies:
            cd    = company_data.get(c, {})
            avg   = cd.get("avg_engagement_rate", 0)
            slope = cd.get("engagement_trend_slope", 0)
            trend = "↑ Growing" if slope > 0.005 else ("↓ Declining" if slope < -0.005 else "→ Stable")
            rows.append([c[:18], f"{avg:.3f}%", trend])

        self._table(slide, headers, rows, Inches(9.1), Inches(1.5), Inches(3.9), Inches(2.4))
        if companies:
            leader = max(companies, key=lambda c: company_data.get(c, {}).get("avg_engagement_rate", 0))
            leader_er = company_data.get(leader, {}).get("avg_engagement_rate", 0.0)
            self._rect(slide, Inches(9.1), Inches(4.15), Inches(3.9), Inches(1.45), C_SOFT_PANEL)
            self._text_box(
                slide,
                self._slide_insight(
                    "slide_07",
                    (
                    f"{leader} has the strongest recent average engagement at {leader_er:.3f}%. "
                    "Use this slide to compare not just the level, but also whether the rolling trend is improving, flat, or fading."
                    ),
                ),
                Inches(9.28), Inches(4.4), Inches(3.55), Inches(0.9),
                font_size=10, color=C_DARK, wrap=True,
            )
        self._footer(slide)

    # ------------------------------------------------------------------
    # SLIDE 8 — Top Performing Videos (white)
    # ------------------------------------------------------------------

    def _slide_08_top_videos(self, slide, companies, company_data, company_colours):
        self._set_bg(slide, C_WHITE)
        self._slide_heading(slide, "Top Performing Videos")

        display_companies = [
            company for company in companies
            if company_data.get(company, {}).get("top_videos")
        ][:4]
        col_w = Inches(6.05)
        card_h = Inches(2.18)
        x_positions = [0.35, 6.93]
        y_positions = [1.32, 4.08]

        for idx, company in enumerate(display_companies):
            x = Inches(x_positions[idx % 2])
            y_start = y_positions[idx // 2]
            hex_colour = company_colours.get(company, MPL_BLUE).lstrip("#")
            rgb = RGBColor(
                int(hex_colour[0:2], 16),
                int(hex_colour[2:4], 16),
                int(hex_colour[4:6], 16),
            )

            # Company header bar
            self._rect(slide, x, Inches(y_start), col_w, Inches(0.42), rgb)
            self._text_box(
                slide, company[:22],
                x, Inches(y_start), col_w, Inches(0.42),
                font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
            )

            videos = company_data.get(company, {}).get("top_videos", [])
            y_vid = y_start + 0.52

            for vid in videos[:1]:
                title   = self._slide_trim(vid.get("title", "N/A"), 60)
                views   = vid.get("view_count", 0)
                likes   = vid.get("like_count", 0)
                cmts    = vid.get("comment_count", 0)
                er      = (likes + cmts) / max(views, 1) * 100
                duration_label = self._duration_label(int(vid.get("duration_seconds", 0)))
                published_label = self._published_label(vid.get("published_at", ""))
                views_per_day = self._video_views_per_day(vid)
                video_url = self._video_url(vid)
                why_it_worked = self._video_performance_hypothesis(vid)
                next_test = self._video_actionable_takeaway(vid)

                self._rect(slide, x, Inches(y_vid), col_w, card_h, C_LIGHT_BG, rounded=True)

                title_shape = self._text_box(
                    slide, title,
                    x + Inches(0.18), Inches(y_vid + 0.1), col_w - Inches(0.36), Inches(0.42),
                    font_size=11, bold=True, color=C_DARK, wrap=True, fit=True, min_font_size=10,
                )
                self._text_box(
                    slide,
                    f"Hero metric: {views_per_day or f'{views:,} lifetime views'}",
                    x + Inches(0.18), Inches(y_vid + 0.54), col_w - Inches(0.36), Inches(0.18),
                    font_size=10, color=C_BLUE, bold=True, fit=True, min_font_size=9,
                )
                meta = " • ".join(part for part in [published_label, duration_label, views_per_day] if part)
                self._text_box(
                    slide, f"Views {views:,} • Likes {likes:,} • ER {self._format_engagement_rate(er)}",
                    x + Inches(0.18), Inches(y_vid + 0.8), col_w - Inches(0.36), Inches(0.14),
                    font_size=8, color=C_GREY, wrap=True, fit=True, min_font_size=8,
                )
                self._text_box(
                    slide, meta,
                    x + Inches(0.18), Inches(y_vid + 0.98), col_w - Inches(0.36), Inches(0.14),
                    font_size=8, color=C_GREY, wrap=True, fit=True, min_font_size=7,
                )
                self._text_box(
                    slide, f"Why it likely worked: {why_it_worked}",
                    x + Inches(0.18), Inches(y_vid + 1.2), col_w - Inches(0.36), Inches(0.28),
                    font_size=8, color=C_MUTED_TEAL, wrap=True, fit=True, min_font_size=8,
                )
                self._text_box(
                    slide, f"Next move: {next_test}",
                    x + Inches(0.18), Inches(y_vid + 1.56), col_w - Inches(0.36), Inches(0.26),
                    font_size=8, color=C_DARK, wrap=True, fit=True, min_font_size=7,
                )
                if video_url:
                    link_shape = self._text_box(
                        slide, "Watch video",
                        x + col_w - Inches(1.1), Inches(y_vid + 1.96), Inches(0.92), Inches(0.14),
                        font_size=8, color=C_BLUE, align=PP_ALIGN.RIGHT,
                    )
                    self._set_textbox_hyperlink(link_shape, video_url, underline=True)

        self._text_box(
            slide,
            self._slide_insight(
                "slide_08",
                self._top_video_slide_summary(companies, company_data),
            ),
            Inches(0.45), Inches(6.95), Inches(12.15), Inches(0.26),
            font_size=9, color=C_GREY, wrap=True, fit=True, min_font_size=8,
        )

        self._footer(slide)

    # ------------------------------------------------------------------
    # SLIDE 9 — Content Topics & Themes (white)
    # ------------------------------------------------------------------

    def _slide_09_topics(self, slide, companies, topic_clusters, company_colours):
        self._set_bg(slide, C_WHITE)
        self._slide_heading(slide, "Content Topics & Themes")

        clusters        = topic_clusters.get("clusters", [])
        company_cov     = topic_clusters.get("company_coverage", {})
        company_themes  = topic_clusters.get("company_theme_labels", {})
        gap_topics      = topic_clusters.get("gap_topics", [])

        if clusters and companies:
            counts = [len(company_cov.get(company, [])) for company in companies]
            colours = [company_colours.get(company, MPL_BLUE) for company in companies]
            chart_img = self._bar_chart_v(
                labels=companies,
                values=counts,
                colours=colours,
                title="Unique Topic Clusters Covered",
                ylabel="Cluster count",
                figsize=(5.8, 3.2),
            )
            slide.shapes.add_picture(chart_img, Inches(0.45), Inches(1.2), Inches(5.6), Inches(3.45))

        self._rect(slide, Inches(6.25), Inches(1.2), Inches(3.0), Inches(4.95), C_LIGHT_BG)
        self._text_box(
            slide, "Representative themes",
            Inches(6.45), Inches(1.38), Inches(2.4), Inches(0.35),
            font_size=12, bold=True, color=C_MUTED_TEAL,
        )
        theme_y = 1.82
        for company in companies[:5]:
            themes = [
                self._humanize_topic_label(label, 2)
                for label in company_themes.get(company, [])[:2]
                if label
            ]
            theme_line = " | ".join(themes) if themes else "No clear recurring theme in the sample"
            self._text_box(
                slide, company,
                Inches(6.45), Inches(theme_y), Inches(2.2), Inches(0.25),
                font_size=10, bold=True, color=C_DARK,
            )
            self._text_box(
                slide, self._slide_trim(theme_line, 54),
                Inches(6.45), Inches(theme_y + 0.22), Inches(2.45), Inches(0.4),
                font_size=10, color=C_DARK, wrap=True,
            )
            theme_y += 0.62

        self._rect(slide, Inches(9.35), Inches(1.2), Inches(3.55), Inches(4.95), C_SOFT_PANEL)
        self._text_box(
            slide, "Strategic whitespace",
            Inches(9.55), Inches(1.38), Inches(2.8), Inches(0.35),
            font_size=12, bold=True, color=C_MUTED_TEAL, align=PP_ALIGN.CENTER,
        )
        gap_text = (
            "\n".join(
                f"• {self._humanize_topic_label(g, 3)}" for g in gap_topics[:6]
            )
            if gap_topics
            else "• Strengthen proof content around the themes competitors treat only as awareness content."
        )
        self._text_box(
            slide, gap_text,
            Inches(9.55), Inches(1.85), Inches(2.95), Inches(3.55),
            font_size=10, color=C_DARK, wrap=True, fit=True, min_font_size=8,
        )
        self._rect(slide, Inches(6.25), Inches(6.25), Inches(6.65), Inches(0.62), C_SOFT_PANEL, rounded=True)
        interpretation = self._slide_insight(
            "slide_09",
            "This slide shows what the market is talking about, not automatically what is working best. The value comes from spotting which themes are overused, underused, or missing proof content the client could own more clearly.",
        )
        self._text_box(
            slide, interpretation,
            Inches(6.45), Inches(6.39), Inches(6.2), Inches(0.38),
            font_size=9, color=C_DARK, wrap=True, fit=True, min_font_size=8,
        )

        self._footer(slide)

    # ------------------------------------------------------------------
    # SLIDE 10 — Video Length Strategy (white)
    # ------------------------------------------------------------------

    def _slide_10_length_strategy(self, slide, companies, company_data, company_colours):
        self._set_bg(slide, C_WHITE)
        self._slide_heading(slide, "Content Format Performance")

        buckets = ["short", "medium", "long"]
        bucket_labels = ["Short (<5 min)", "Medium (5–15 min)", "Long (>15 min)"]
        bucket_colours = [MPL_TEAL, MPL_BLUE, MPL_AMBER]

        # Grouped bar chart
        fig, ax = plt.subplots(figsize=(8, 3.5))
        ax.set_facecolor("white")
        fig.patch.set_facecolor("white")
        for spine in ax.spines.values():
            spine.set_visible(False)
        ax.yaxis.grid(True, color="#E2E8F0", linewidth=0.8)
        ax.set_axisbelow(True)

        n_companies = len(companies)
        n_buckets   = len(buckets)
        bar_width   = 0.6 / n_buckets
        x = np.arange(n_companies)

        for bi, (bucket, label, colour) in enumerate(zip(buckets, bucket_labels, bucket_colours)):
            ers = [
                company_data.get(c, {}).get(bucket, {}).get("avg_engagement_rate", 0)
                for c in companies
            ]
            offset = (bi - n_buckets / 2 + 0.5) * bar_width
            bars = ax.bar(x + offset, ers, width=bar_width * 0.9, label=label, color=colour, alpha=0.88)
            for bar, val in zip(bars, ers):
                if val > 0:
                    ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.002,
                            f"{val:.2f}", ha="center", va="bottom", fontsize=7, color="#1E293B")

        ax.set_xticks(x)
        ax.set_xticklabels([c[:12] for c in companies], fontsize=9, color="#64748B")
        ax.set_ylabel("Avg Engagement Rate (%)", fontsize=9, color="#64748B")
        ax.set_title("Engagement Rate by Video Length Bucket", fontsize=10, color="#1E293B", pad=8)
        ax.legend(fontsize=9, frameon=False)
        plt.tight_layout()
        chart_img = self._fig_to_image(fig)
        slide.shapes.add_picture(chart_img, Inches(0.4), Inches(1.2), Inches(8.5), Inches(3.8))

        headers = ["Company", "Best", "Confidence", "Next test"]
        rows = []
        for c in companies:
            cd   = company_data.get(c, {})
            best = self._length_label(cd.get("best_performing_length", "—"))
            confidence = cd.get("length_confidence", "LOW").title()
            rec  = self._length_recommendation(cd)
            rows.append([self._company_short_label(c, 12), best, confidence, rec])

        self._table(slide, headers, rows, Inches(8.86), Inches(1.42), Inches(4.09), Inches(2.55))
        self._rect(slide, Inches(0.38), Inches(5.28), Inches(12.45), Inches(0.72), C_SOFT_PANEL, rounded=True)
        self._text_box(
            slide,
            self._slide_insight(
                "slide_10",
                self._length_strategy_takeaway(companies, company_data),
            ),
            Inches(0.58), Inches(5.44), Inches(12.0), Inches(0.3),
            font_size=9, color=C_DARK, wrap=True, fit=True, min_font_size=8,
        )
        self._footer(slide)

    # ------------------------------------------------------------------
    # SLIDE 11 — Video SEO Analysis (white)
    # ------------------------------------------------------------------

    def _slide_11_seo(self, slide, companies, seo_scores, company_data, company_colours):
        self._set_bg(slide, C_WHITE)
        self._slide_heading(slide, "Discovery & Search Visibility")

        dimensions = ["Title Length", "Description\nDepth", "Timestamps", "Tags", "Keywords"]
        dim_keys   = ["title_length_score", "description_depth", "has_timestamps_pct",
                      "tags_count_avg", "keyword_in_title_score"]

        # Normalise tags to 0-100
        max_tags = max(
            (seo_scores.get(c, {}).get("breakdown", {}).get("tags_count_avg", 0) for c in companies),
            default=10,
        )
        max_tags = max(max_tags, 1)

        fig, ax = plt.subplots(figsize=(5.5, 4.0), subplot_kw={"polar": True})
        fig.patch.set_facecolor("white")
        ax.set_facecolor("white")

        num_vars = len(dimensions)
        angles = [n / float(num_vars) * 2 * math.pi for n in range(num_vars)]
        angles += angles[:1]

        ax.set_xticks(angles[:-1])
        ax.set_xticklabels(dimensions, size=8, color="#1E293B")
        ax.set_ylim(0, 100)
        ax.set_yticks([25, 50, 75, 100])
        ax.set_yticklabels(["25", "50", "75", "100"], size=7, color="#94A3B8")
        ax.grid(color="#E2E8F0", linewidth=0.7)

        for c in companies:
            bd   = seo_scores.get(c, {}).get("breakdown", {})
            vals = [
                bd.get("title_length_score", 0),
                bd.get("description_depth", 0),
                bd.get("has_timestamps_pct", 0),
                min(bd.get("tags_count_avg", 0) / max_tags * 100, 100),
                bd.get("keyword_in_title_score", 0),
            ]
            vals += vals[:1]
            hex_c = company_colours.get(c, MPL_BLUE)
            ax.plot(angles, vals, color=hex_c, linewidth=2, label=c)
            ax.fill(angles, vals, color=hex_c, alpha=0.12)

        ax.legend(loc="lower right", bbox_to_anchor=(1.35, -0.12), fontsize=8, frameon=False)
        plt.tight_layout()
        radar_img = self._fig_to_image(fig)
        slide.shapes.add_picture(radar_img, Inches(0.3), Inches(1.2), Inches(5.8), Inches(4.2))

        # SEO table (right side)
        headers = ["Company", "SEO Score", "Title score", "Description", "Timestamps"]
        rows = []
        for c in companies:
            sd = seo_scores.get(c, {})
            bd = sd.get("breakdown", {})
            rows.append([
                c[:14],
                f"{sd.get('seo_score', 0):.1f}",
                f"{bd.get('title_length_score', 0):.0f}",
                f"{bd.get('description_depth', 0):.0f}",
                f"{bd.get('has_timestamps_pct', 0):.0f}%",
            ])

        self._table(slide, headers, rows, Inches(6.4), Inches(1.5), Inches(6.6), Inches(2.8))
        timestamp_zero_count = sum(
            1 for company in companies
            if seo_scores.get(company, {}).get("breakdown", {}).get("has_timestamps_pct", 0.0) == 0.0
        )
        explanation = self._slide_insight(
            "slide_11",
            self._seo_takeaway(companies, seo_scores),
        )
        self._text_box(
            slide, explanation,
            Inches(6.45), Inches(4.65), Inches(6.3), Inches(1.1),
            font_size=10, color=C_GREY, wrap=True, fit=True, min_font_size=8,
        )
        self._footer(slide)

    # ------------------------------------------------------------------
    # SLIDE 12 — Gap Analysis (light)
    # ------------------------------------------------------------------

    def _slide_12_gap_analysis(self, slide, gap_analysis: str, opp_scores: list):
        self._set_bg(slide, C_PAPER)
        self._slide_heading(slide, "Strategic Whitespace Opportunities", dark=False)

        if opp_scores:
            labels = [self._humanize_topic_label(o.get("topic", ""), 2)[:28] for o in opp_scores[:5]]
            values = [o.get("opportunity_score", 0) for o in opp_scores[:5]]
            colours = [COMPANY_COLOURS[i % len(COMPANY_COLOURS)] for i in range(len(values))]
            chart_img = self._bar_chart_h(
                labels=labels,
                values=values,
                colours=colours,
                title="Opportunity ranking (higher = stronger whitespace signal)",
                xlabel="Opportunity score",
                xlim=(0, 100),
                figsize=(7.0, 3.6),
            )
            slide.shapes.add_picture(chart_img, Inches(0.35), Inches(1.2), Inches(7.4), Inches(3.95))

        top3 = opp_scores[:3]
        self._text_box(
            slide, "Where To Test Next",
            Inches(8.28), Inches(1.18), Inches(4.6), Inches(0.28),
            font_size=13, bold=True, color=C_MUTED_TEAL,
        )
        y_pos = 1.58
        for rank, opp in enumerate(top3, 1):
            self._rect(slide, Inches(8.1), Inches(y_pos), Inches(4.82), Inches(1.36), C_SOFT_PANEL, rounded=True)
            topic = self._humanize_topic_label(opp.get("topic", "—"), 2)[:30]
            score = opp.get("opportunity_score", 0)
            evidence = self._slide_trim(
                opp.get("supporting_evidence", f"Trend {opp.get('trend_interest', 0):.0f}/100; scarcity {opp.get('scarcity_score', 0):.0f}/100."),
                62,
            )
            why = self._slide_trim(opp.get("opportunity_brief", opp.get("recommendation_type", "")), 50)
            experiment = self._slide_trim(opp.get("suggested_experiment", "Test one tightly scoped educational video on this theme."), 48)
            signal = self._slide_trim(opp.get("signal_to_watch", "Track first-14-day views-per-day and engagement rate."), 46)
            self._text_box(
                slide, f"{rank}. {topic} · {score:.0f}/100",
                Inches(8.28), Inches(y_pos + 0.08), Inches(4.3), Inches(0.18),
                font_size=10, bold=True, color=C_DARK, fit=True, min_font_size=9,
            )
            self._text_box(
                slide, f"Why: {why}",
                Inches(8.28), Inches(y_pos + 0.32), Inches(4.35), Inches(0.16),
                font_size=8, color=C_MUTED_TEAL, wrap=True, fit=True, min_font_size=7,
            )
            self._text_box(
                slide, f"Evidence: {evidence}",
                Inches(8.28), Inches(y_pos + 0.52), Inches(4.35), Inches(0.14),
                font_size=8, color=C_GREY, wrap=True, fit=True, min_font_size=7,
            )
            self._text_box(
                slide, f"Test: {experiment}",
                Inches(8.28), Inches(y_pos + 0.74), Inches(4.35), Inches(0.14),
                font_size=8, color=C_DARK, wrap=True, fit=True, min_font_size=7,
            )
            self._text_box(
                slide, f"Watch: {signal}",
                Inches(8.28), Inches(y_pos + 0.96), Inches(4.35), Inches(0.14),
                font_size=8, color=C_DARK, wrap=True, fit=True, min_font_size=7,
            )
            y_pos += 1.46

        excerpt = self._slide_insight("slide_12", self._gap_excerpt(gap_analysis, opp_scores))
        self._text_box(
            slide, excerpt,
            Inches(0.42), Inches(5.4), Inches(7.3), Inches(1.2),
            font_size=10, color=C_DARK, wrap=True, fit=True, min_font_size=9,
        )
        self._footer(slide, dark=False)

    # ------------------------------------------------------------------
    # SLIDE 13 — Recommendations (white)
    # ------------------------------------------------------------------

    def _slide_13_priority_moves(self, slide, recommendations: list, priority_filter: str):
        self._set_bg(slide, C_WHITE)
        is_high = priority_filter.lower() == "high"
        slide_key = "slide_13" if is_high else "slide_14"
        heading = "Priority Growth Moves — High Priority" if is_high else "Priority Growth Moves — Medium Priority"
        filtered = [rec for rec in recommendations if str(rec.get("priority", "")).lower() == priority_filter.lower()]
        if not filtered:
            filtered = recommendations[:2] if is_high else recommendations[2:5]

        self._slide_heading(slide, heading)
        intro = self._slide_trim(self._slide_insight(
            slide_key,
            (
                "These are the highest-priority actions because they appear most likely to improve client growth leverage quickly."
                if is_high else
                "These medium-priority actions strengthen execution quality, reduce blind spots, and improve learning speed after the highest-priority moves are underway."
            ),
        ), 180 if is_high else 170)
        self._text_box(
            slide, intro,
            Inches(0.45), Inches(1.03), Inches(12.1), Inches(0.3),
            font_size=10, color=C_GREY, wrap=True, fit=True, min_font_size=9,
        )

        y_pos = 1.35
        max_cards = 2 if is_high else 3
        for i, rec in enumerate(filtered[:max_cards]):
            title   = self._slide_trim(rec.get("title", f"Recommendation {i+1}"), 58)
            action  = self._slide_trim(self._clean_markdown(rec.get("action", "")), 88)
            rationale = self._slide_trim(self._clean_markdown(rec.get("why_this_matters") or rec.get("data_rationale", "")), 72)
            impact  = self._slide_trim(self._clean_markdown(rec.get("business_impact") or rec.get("expected_impact", "")), 58)
            kpi = self._slide_trim(self._clean_markdown(rec.get("success_kpi", "")), 44)
            evidence = self._slide_trim(self._clean_markdown(rec.get("supporting_evidence") or rec.get("data_rationale", "")), 64)
            priority = rec.get("priority", "medium").lower()
            prio_colour = C_RED if priority == "high" else C_AMBER

            card_height = Inches(2.0 if is_high else 1.7)
            self._rect(slide, Inches(0.35), Inches(y_pos), Inches(12.15), card_height, C_LIGHT_BG, rounded=True)
            self._rect(slide, Inches(0.45), Inches(y_pos + 0.12), Inches(0.42), Inches(0.42), C_MUTED_TEAL, rounded=True)
            self._text_box(
                slide, str(i + 1),
                Inches(0.45), Inches(y_pos + 0.12), Inches(0.42), Inches(0.42),
                font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
            )

            self._text_box(
                slide, title,
                Inches(1.0), Inches(y_pos + 0.12), Inches(8.4), Inches(0.24),
                font_size=13, bold=True, color=C_DARK, fit=True, min_font_size=11,
            )

            self._text_box(
                slide, priority.upper(),
                Inches(10.85), Inches(y_pos + 0.12), Inches(1.0), Inches(0.24),
                font_size=10, bold=True, color=prio_colour, align=PP_ALIGN.CENTER,
            )

            self._text_box(
                slide, action,
                Inches(1.0), Inches(y_pos + 0.44), Inches(8.65), Inches(0.32),
                font_size=12 if is_high else 11, color=C_DARK, wrap=True, fit=True, min_font_size=9,
            )

            if rationale:
                self._text_box(
                    slide, "Why this matters: " + rationale,
                    Inches(1.0), Inches(y_pos + 0.84), Inches(8.65), Inches(0.2),
                    font_size=10 if is_high else 9, color=C_GREY, wrap=True, fit=True, min_font_size=8,
                )
            if impact:
                self._text_box(
                    slide, "Business impact: " + impact,
                    Inches(1.0), Inches(y_pos + 1.12), Inches(5.45), Inches(0.18),
                    font_size=10 if is_high else 9, color=C_MUTED_TEAL, italic=True, wrap=True, fit=True, min_font_size=8,
                )
            if kpi:
                self._rect(slide, Inches(9.45), Inches(y_pos + 0.58), Inches(2.32), Inches(0.52), C_SOFT_PANEL, rounded=True)
                self._text_box(
                    slide, "KPI: " + kpi,
                    Inches(9.62), Inches(y_pos + 0.74), Inches(1.98), Inches(0.18),
                    font_size=9, color=C_DARK, wrap=True, fit=True, min_font_size=8, align=PP_ALIGN.CENTER,
                )
            if evidence:
                self._text_box(
                    slide, "Evidence: " + evidence,
                    Inches(1.0), Inches(y_pos + (1.42 if is_high else 1.28)), Inches(10.9), Inches(0.16),
                    font_size=8, color=C_GREY, wrap=True, fit=True, min_font_size=7,
                )

            y_pos += 2.12 if is_high else 1.84

        self._footer(slide)

    # ------------------------------------------------------------------
    # SLIDE 14 — Summary Scorecard (light)
    # ------------------------------------------------------------------

    def _slide_14_scorecard(self, slide, companies, company_data, rpi_scores, seo_scores, report_date, target_company: str):
        self._set_bg(slide, C_PAPER)
        self._slide_heading(slide, "Company Growth Scorecard", dark=False)

        headers = ["Company", "RPI Rank", "Engagement %", "Consistency", "Funnel", "SEO Score", "Grade"]
        rows = []
        for c in companies:
            cd   = company_data.get(c, {})
            rpi  = rpi_scores.get(c, {})
            seo  = seo_scores.get(c, {})

            rank      = rpi.get("rank", "—")
            eng       = f"{cd.get('avg_engagement_rate', 0):.3f}%"
            cons      = f"{cd.get('consistency_score', 0):.1f}"
            funnel    = cd.get("funnel_label", "—")
            seo_s     = f"{seo.get('seo_score', 0):.1f}"
            grade     = self._rpi_grade(rpi.get("rpi_score", 0))
            rows.append([c[:18], str(rank), eng, cons, funnel, seo_s, grade])

        self._table(
            slide, headers, rows,
            Inches(0.4), Inches(1.4), Inches(12.5), Inches(4.2),
            dark=False,
        )

        # Footer credit
        self._text_box(
            slide,
            f"Analysis generated {report_date} | Data source: YouTube Data API v3",
            Inches(0.4), Inches(6.38), Inches(12.5), Inches(0.26),
            font_size=10, color=C_GREY, align=PP_ALIGN.CENTER,
        )
        self._text_box(
            slide,
            "Grade bands: A = 75+, B = 50–74.9, C = 25–49.9, D = below 25 on the RPI scale.",
            Inches(0.45), Inches(6.63), Inches(12.3), Inches(0.24),
            font_size=10, color=C_GREY, align=PP_ALIGN.CENTER, wrap=True,
        )
        if companies:
            ranked = sorted(companies, key=lambda c: rpi_scores.get(c, {}).get("rank", 999))
            leader = ranked[0]
            closing = self._slide_trim(self._slide_insight(
                "slide_15",
                self._scorecard_closing(companies, company_data, rpi_scores, seo_scores, target_company),
            ), 190)
            self._rect(slide, Inches(0.6), Inches(6.76), Inches(12.0), Inches(0.34), C_SOFT_PANEL, rounded=True)
            self._text_box(
                slide, closing,
                Inches(0.78), Inches(6.83), Inches(11.65), Inches(0.22),
                font_size=9, color=C_DARK, align=PP_ALIGN.CENTER, wrap=True, fit=True, min_font_size=8,
            )
        self._footer(slide, dark=False)

    def _slide_15_action_plan(self, slide, action_plan: str, recommendations: list[dict], target_company: str, target_data: dict, opp_scores: list[dict]):
        self._set_bg(slide, C_WHITE)
        self._slide_heading(slide, "90-Day Action Plan")

        intro = self._slide_insight(
            "slide_16",
            self._action_plan_takeaway(target_company, target_data, recommendations, opp_scores),
        )
        self._text_box(
            slide, intro,
            Inches(0.45), Inches(1.02), Inches(12.0), Inches(0.3),
            font_size=10, color=C_GREY, wrap=True, fit=True, min_font_size=9,
        )

        phases = self._action_plan_sections(action_plan)
        headers = ["Week 1–2", "Month 1", "Month 2–3"]
        bodies = self._action_plan_card_bodies(phases, recommendations, target_company, target_data, opp_scores)
        x_positions = [Inches(0.45), Inches(4.48), Inches(8.51)]
        for idx, (header, body) in enumerate(zip(headers, bodies)):
            x = x_positions[idx]
            self._rect(slide, x, Inches(1.55), Inches(3.7), Inches(4.55), C_LIGHT_BG, rounded=True)
            self._text_box(
                slide, header,
                x + Inches(0.18), Inches(1.78), Inches(2.5), Inches(0.3),
                font_size=14, bold=True, color=C_BLUE,
            )
            self._text_box(
                slide, body,
                x + Inches(0.18), Inches(2.15), Inches(3.35), Inches(2.95),
                font_size=12, color=C_DARK, wrap=True, fit=True, min_font_size=9,
            )

        top_priorities = ", ".join(
            self._slide_trim(rec.get("title", ""), 42)
            for rec in recommendations[:2]
            if rec.get("title")
        )
        self._rect(slide, Inches(0.45), Inches(6.38), Inches(11.95), Inches(0.42), C_SOFT_PANEL, rounded=True)
        self._text_box(
            slide,
            f"Priority focus for the first month: {top_priorities or 'Use the highest-priority recommendation as the opening sprint.'}",
            Inches(0.62), Inches(6.5), Inches(11.6), Inches(0.18),
            font_size=10, color=C_DARK, wrap=True, fit=True, min_font_size=8,
        )
        self._footer(slide)

    # ==================================================================
    # Chart helpers
    # ==================================================================

    def _bar_chart_h(
        self,
        labels: list[str],
        values: list[float],
        colours: list[str],
        title: str = "",
        xlabel: str = "",
        xlim: tuple | None = None,
        figsize: tuple = (8, 3.5),
    ) -> io.BytesIO:
        """Horizontal bar chart. Returns BytesIO buffer."""
        fig, ax = plt.subplots(figsize=figsize)
        ax.set_facecolor("white")
        fig.patch.set_facecolor("white")
        for spine in ax.spines.values():
            spine.set_visible(False)
        ax.xaxis.grid(True, color="#E2E8F0", linewidth=0.8)
        ax.set_axisbelow(True)

        y_pos = range(len(labels))
        hex_colours = [c.lstrip("#") for c in colours]
        rgb_colours = [
            (int(h[0:2],16)/255, int(h[2:4],16)/255, int(h[4:6],16)/255)
            for h in hex_colours
        ]
        bars = ax.barh(list(y_pos), values, color=rgb_colours, height=0.55, alpha=0.90)
        ax.set_yticks(list(y_pos))
        ax.set_yticklabels([lb[:18] for lb in labels], fontsize=9, color="#1E293B")
        ax.tick_params(axis="x", colors="#64748B")
        if xlabel:
            ax.set_xlabel(xlabel, fontsize=9, color="#64748B")
        if title:
            ax.set_title(title, fontsize=10, color="#1E293B", pad=6)
        if xlim:
            ax.set_xlim(*xlim)

        for bar, val in zip(bars, values):
            ax.text(
                bar.get_width() + max(values) * 0.01,
                bar.get_y() + bar.get_height() / 2,
                f"{val:,.0f}",
                va="center", fontsize=8, color="#1E293B",
            )

        plt.tight_layout()
        return self._fig_to_image(fig)

    def _bar_chart_v(
        self,
        labels: list[str],
        values: list[float],
        colours: list[str],
        title: str = "",
        ylabel: str = "",
        ylim: tuple | None = None,
        figsize: tuple = (6, 3.5),
    ) -> io.BytesIO:
        """Vertical bar chart. Returns BytesIO buffer."""
        fig, ax = plt.subplots(figsize=figsize)
        ax.set_facecolor("white")
        fig.patch.set_facecolor("white")
        for spine in ax.spines.values():
            spine.set_visible(False)
        ax.yaxis.grid(True, color="#E2E8F0", linewidth=0.8)
        ax.set_axisbelow(True)

        hex_colours = [c.lstrip("#") for c in colours]
        rgb_colours = [
            (int(h[0:2],16)/255, int(h[2:4],16)/255, int(h[4:6],16)/255)
            for h in hex_colours
        ]
        bars = ax.bar(range(len(labels)), values, color=rgb_colours, width=0.55, alpha=0.90)
        ax.set_xticks(range(len(labels)))
        ax.set_xticklabels([lb[:12] for lb in labels], fontsize=9, color="#64748B")
        ax.tick_params(axis="y", colors="#64748B")
        if ylabel:
            ax.set_ylabel(ylabel, fontsize=9, color="#64748B")
        if title:
            ax.set_title(title, fontsize=9, color="#1E293B", pad=6)
        if ylim:
            ax.set_ylim(*ylim)

        for bar, val in zip(bars, values):
            ax.text(
                bar.get_x() + bar.get_width() / 2,
                bar.get_height() + max(values) * 0.01,
                f"{val:.1f}",
                ha="center", fontsize=8, color="#1E293B",
            )

        plt.tight_layout()
        return self._fig_to_image(fig)

    def _stacked_bar_chart(
        self,
        companies: list[str],
        series: dict[str, list[float]],
        colours: list[str],
        title: str = "",
    ) -> io.BytesIO:
        """Horizontal stacked bar chart (for funnel distribution)."""
        fig, ax = plt.subplots(figsize=(8, 3.5))
        ax.set_facecolor("white")
        fig.patch.set_facecolor("white")
        for spine in ax.spines.values():
            spine.set_visible(False)
        ax.xaxis.grid(True, color="#E2E8F0", linewidth=0.8)
        ax.set_axisbelow(True)

        y_pos  = range(len(companies))
        lefts  = [0.0] * len(companies)
        hex_cs = [c.lstrip("#") for c in colours]
        rgb_cs = [(int(h[0:2],16)/255, int(h[2:4],16)/255, int(h[4:6],16)/255) for h in hex_cs]

        for (label, vals), colour in zip(series.items(), rgb_cs):
            ax.barh(list(y_pos), vals, left=lefts, color=colour, height=0.55,
                    label=label, alpha=0.90)
            lefts = [l + v for l, v in zip(lefts, vals)]

        ax.set_yticks(list(y_pos))
        ax.set_yticklabels([c[:18] for c in companies], fontsize=9, color="#1E293B")
        ax.set_xlabel("Percentage (%)", fontsize=9, color="#64748B")
        ax.set_xlim(0, 105)
        if title:
            ax.set_title(title, fontsize=10, color="#1E293B", pad=6)
        ax.set_xlabel("")
        ax.legend(
            fontsize=7,
            frameon=False,
            loc="upper center",
            bbox_to_anchor=(0.5, -0.12),
            ncol=2,
        )
        fig.subplots_adjust(bottom=0.24)
        plt.tight_layout()
        return self._fig_to_image(fig)

    @staticmethod
    def _fig_to_image(fig) -> io.BytesIO:
        """Save matplotlib figure to BytesIO and close it."""
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=180, bbox_inches="tight")
        plt.close(fig)
        buf.seek(0)
        return buf

    # ==================================================================
    # python-pptx helper methods
    # ==================================================================

    @staticmethod
    def _set_bg(slide, colour: RGBColor):
        """Set slide solid background colour."""
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = colour

    def _slide_heading(self, slide, text: str, dark: bool = False):
        """Add a styled slide heading."""
        colour = C_WHITE if dark else C_DARK
        self._text_box(
            slide, text,
            Inches(0.4), Inches(0.18), Inches(12.0), Inches(0.85),
            font_size=28, bold=True, color=colour, font_face=FONT_HEADING,
        )
        # Subtle underline rule
        rule_colour = C_TEAL if dark else C_BLUE
        self._rect(slide, Inches(0.4), Inches(0.98), Inches(12.5), Inches(0.03), rule_colour)

    @staticmethod
    def _text_box(
        slide,
        text: str,
        left, top, width, height,
        font_size: int = 14,
        bold: bool = False,
        italic: bool = False,
        color: RGBColor = C_DARK,
        align: PP_ALIGN = PP_ALIGN.LEFT,
        font_face: str = FONT_BODY,
        wrap: bool = False,
        fit: bool = False,
        min_font_size: int = 8,
    ):
        """Add a text box with consistent styling."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = wrap
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size   = Pt(font_size)
        run.font.bold   = bold
        run.font.italic = italic
        run.font.name   = font_face
        run.font.color.rgb = color
        if fit:
            try:
                tf.fit_text(font_family=font_face, max_size=font_size)
            except Exception:
                pass
            for paragraph in tf.paragraphs:
                for paragraph_run in paragraph.runs:
                    if paragraph_run.font.size and paragraph_run.font.size.pt < min_font_size:
                        paragraph_run.font.size = Pt(min_font_size)
        return txBox

    @staticmethod
    def _rect(slide, left, top, width, height, colour: RGBColor, rounded: bool = False):
        """Add a filled rectangle (no border)."""
        shape_type = 5 if not rounded else 6  # MSO_SHAPE_TYPE constants via index
        from pptx.util import Emu
        shape = slide.shapes.add_shape(
            1 if not rounded else 5,  # MSO_AUTO_SHAPE_TYPE: RECTANGLE=1, ROUNDED_RECT=5
            left, top, width, height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colour
        shape.line.fill.background()

    def _table(
        self,
        slide,
        headers: list[str],
        rows: list[list[str]],
        left, top, width, height,
        dark: bool = False,
    ):
        """Add a styled table with a header row."""
        n_cols = len(headers)
        n_rows = len(rows) + 1

        tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
        col_w = width // n_cols
        for i in range(n_cols):
            tbl.columns[i].width = col_w

        # Header row
        for ci, header in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_BLUE
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Emu(45720)
            tf.margin_right = Emu(45720)
            tf.margin_top = Emu(22860)
            tf.margin_bottom = Emu(22860)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = tf.paragraphs[0].add_run()
            run.text = header
            run.font.bold  = True
            run.font.size  = Pt(12)
            run.font.name  = FONT_HEADING
            run.font.color.rgb = C_WHITE

        # Data rows
        for ri, row in enumerate(rows):
            for ci, value in enumerate(row):
                cell = tbl.cell(ri + 1, ci)
                fill_colour = C_LIGHT_BG if (ri % 2 == 0 and not dark) else (C_NAVY if dark else C_WHITE)
                if dark:
                    fill_colour = RGBColor(22, 44, 66) if ri % 2 == 0 else RGBColor(13, 27, 42)
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill_colour
                tf = cell.text_frame
                tf.word_wrap = True
                tf.margin_left = Emu(45720)
                tf.margin_right = Emu(45720)
                tf.margin_top = Emu(22860)
                tf.margin_bottom = Emu(22860)
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                run = tf.paragraphs[0].add_run()
                run.text = str(value)
                run.font.size  = Pt(12)
                run.font.name  = FONT_BODY
                run.font.color.rgb = C_WHITE if dark else C_DARK

    @staticmethod
    def _set_textbox_hyperlink(shape, url: str, underline: bool = False):
        if not url:
            return
        try:
            shape.click_action.hyperlink.address = url
        except Exception:
            pass
        try:
            paragraph = shape.text_frame.paragraphs[0]
            if paragraph.runs:
                paragraph.runs[0].hyperlink.address = url
                if underline:
                    paragraph.runs[0].font.underline = True
        except Exception:
            pass

    @staticmethod
    def _footer(slide, dark: bool = False):
        """Add bottom-right footer watermark."""
        colour = C_GREY if not dark else RGBColor(148, 163, 184)
        txBox = slide.shapes.add_textbox(
            Inches(9.2), Inches(7.15), Inches(4.0), Inches(0.28)
        )
        tf = txBox.text_frame
        p  = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = FOOTER_TEXT
        run.font.size  = Pt(9)
        run.font.name  = FONT_BODY
        run.font.color.rgb = colour

    # ==================================================================
    # Utility helpers
    # ==================================================================

    @staticmethod
    def _channel_age_years(published_at: str) -> str:
        """Return channel age in years from ISO 8601 date string."""
        if not published_at:
            return "—"
        try:
            dt  = datetime.fromisoformat(published_at.replace("Z", "+00:00"))
            age = (datetime.now(dt.tzinfo) - dt).days / 365.25
            return f"{age:.1f}"
        except (ValueError, AttributeError):
            return "—"

    @staticmethod
    def _cadence_label(mean_gap_days: float) -> str:
        """Convert mean_gap_days to a human-readable cadence label."""
        if mean_gap_days <= 0:
            return "—"
        if mean_gap_days < 1:
            return "Burst / same-day"
        if mean_gap_days <= 9:
            return "Weekly"
        if mean_gap_days <= 18:
            return "Bi-weekly"
        if mean_gap_days <= 40:
            return "Monthly"
        return "Irregular"

    @staticmethod
    def _format_gap_days(mean_gap_days: float) -> str:
        """Format average upload gap so tiny values don't display misleadingly."""
        if mean_gap_days <= 0:
            return "—"
        if mean_gap_days < 1:
            return "<1 day"
        return f"{mean_gap_days:.1f}"

    @staticmethod
    def _clean_markdown(text: str) -> str:
        cleaned = re.sub(r"\*\*", "", text or "")
        cleaned = re.sub(r"`", "", cleaned)
        cleaned = re.sub(r"\s+", " ", cleaned).strip()
        return cleaned

    def _slide_insight(self, slide_key: str, fallback: str) -> str:
        insight_map = getattr(self, "_active_interpretations", {}) or {}
        candidate = self._clean_markdown(str(insight_map.get(slide_key, "") or ""))
        candidate = candidate.replace("...", ".")
        return candidate or fallback

    def _executive_bullets(self, company_data: dict, rpi_scores: dict, seo_scores: dict) -> list[str]:
        companies = list(rpi_scores.keys())
        if not companies:
            return []
        ranked = sorted(companies, key=lambda c: rpi_scores.get(c, {}).get("rank", 999))
        leader = ranked[0]
        runner_up = ranked[1] if len(ranked) > 1 else None
        best_consistency = max(companies, key=lambda c: company_data.get(c, {}).get("consistency_score", 0))
        best_seo = max(companies, key=lambda c: seo_scores.get(c, {}).get("seo_score", 0))
        bullets = [
            f"{leader} ranks first with an RPI of {rpi_scores.get(leader, {}).get('rpi_score', 0.0):.1f} and recent engagement of {company_data.get(leader, {}).get('avg_engagement_rate', 0.0):.3f}%.",
            (
                f"{best_consistency} has the strongest publishing discipline at "
                f"{company_data.get(best_consistency, {}).get('consistency_score', 0.0):.1f}/100."
            ),
            (
                f"{best_seo} leads on video SEO at {seo_scores.get(best_seo, {}).get('seo_score', 0.0):.1f}/100, "
                "showing the strongest packaging for search and discovery."
            ),
        ]
        if runner_up:
            bullets[0] = (
                f"{leader} leads overall with an RPI of {rpi_scores.get(leader, {}).get('rpi_score', 0.0):.1f}, "
                f"ahead of {runner_up} at {rpi_scores.get(runner_up, {}).get('rpi_score', 0.0):.1f}."
            )
        return [self._slide_trim(bullet, 128) for bullet in bullets]

    @staticmethod
    def _slide_trim(text: str, limit: int) -> str:
        text = " ".join((text or "").split())
        if len(text) <= limit:
            return text
        clipped = text[:limit]
        stop_positions = [clipped.rfind(marker) for marker in (". ", "; ", ", ")]
        stop = max(stop_positions)
        if stop > int(limit * 0.6):
            clipped = clipped[: stop + 1]
        else:
            clipped = clipped.rsplit(" ", 1)[0].strip()
        return clipped.rstrip(" ,;:.") + "."

    @staticmethod
    def _channel_age_years_value(published_at: str) -> float:
        if not published_at:
            return 999.0
        try:
            dt = datetime.fromisoformat(published_at.replace("Z", "+00:00"))
            return round((datetime.now(dt.tzinfo) - dt).days / 365.25, 2)
        except (ValueError, AttributeError):
            return 999.0

    def _summary_sections(self, text: str) -> tuple[str, list[str]]:
        """Convert a long AI summary into a short intro plus compact bullets."""
        cleaned = self._clean_markdown(text)
        cleaned = re.sub(r"Executive Summary:?\s*", "", cleaned, flags=re.IGNORECASE)
        sentences = [
            self._smart_truncate(chunk.strip(), 170)
            for chunk in re.split(r"(?<=[.!?])\s+", cleaned)
            if chunk.strip()
        ]
        if not sentences:
            return "No summary generated.", []
        intro = sentences[0]
        bullets = []
        for sentence in sentences[1:]:
            bullet = self._smart_truncate(sentence, 130)
            if bullet and bullet not in bullets:
                bullets.append(bullet)
            if len(bullets) >= 3:
                break
        return intro, bullets

    def _executive_summary_parts(self, text: str) -> tuple[str, list[str], str]:
        cleaned = self._clean_markdown(text)
        if not cleaned:
            return "No summary generated.", [], "Use the benchmark slides to identify the first growth lever to test."

        intro = ""
        bullets: list[str] = []
        closing = ""
        for raw_line in cleaned.splitlines():
            line = raw_line.strip()
            if not line:
                continue
            if line.lower().startswith("what this means:"):
                closing = line.split(":", 1)[1].strip() or line
            elif line.startswith("- "):
                bullets.append(self._slide_trim(line[2:].strip(), 126))
            elif not intro:
                intro = self._slide_trim(line, 210)
            else:
                bullets.append(self._slide_trim(line, 126))

        if not intro:
            intro, fallback_bullets = self._summary_sections(cleaned)
            bullets = bullets or fallback_bullets
        if not closing:
            ranked_hint = bullets[0] if bullets else intro
            closing = self._slide_trim(
                f"What this means: {ranked_hint} Use the rest of the deck to choose the next action with the clearest business upside.",
                170,
            )
            closing = closing.replace("What this means: What this means:", "What this means:")
        return intro, bullets[:5], closing

    @staticmethod
    def _smart_truncate(text: str, limit: int) -> str:
        text = " ".join((text or "").split())
        if len(text) <= limit:
            return text
        clipped = text[: limit - 1].rsplit(" ", 1)[0].strip()
        return (clipped or text[: limit - 1]).strip() + "…"

    def _humanize_topic_label(self, label: str, max_terms: int = 3) -> str:
        """Turn noisy cluster labels into cleaner client-facing topic names."""
        strategic_labels = {
            "Product Education",
            "AI Automation",
            "Customer Proof & ROI",
            "Enterprise Security & Trust",
            "Integrations & Workflows",
            "Onboarding & Adoption",
            "Product Launches & Updates",
            "Industry POV & Strategy",
            "Comparison & Evaluation",
            "Use Cases by Role",
        }
        if label in strategic_labels:
            return label
        text = self._clean_markdown(label).replace("|", "/").replace("_", " ")
        parts = [part.strip() for part in text.split("/") if part.strip()]
        cleaned_parts: list[str] = []
        seen_roots: set[str] = set()
        noise_tokens = {
            "http", "https", "www", "com", "video", "videos", "youtube",
            "channel", "official", "utm", "medium", "source", "campaign",
            "mailchimp", "hubspot", "salesforce", "mypromovideos", "clickhubspot",
            "intuit",
        }

        for part in parts:
            part = part.replace("ai foundry", "AI foundry")
            part = part.replace("agentforce", "Agentforce")
            words = []
            for raw_word in re.findall(r"[A-Za-z0-9\+\-]+", part):
                lower = raw_word.lower()
                root = lower.rstrip("s")
                if lower in noise_tokens or root in noise_tokens:
                    continue
                if len(lower) <= 1:
                    continue
                if lower.isdigit():
                    continue
                if len(lower) > 18:
                    continue
                if len(set(lower)) == 1 and len(lower) > 3:
                    continue
                if lower.endswith("edium"):
                    continue
                words.append(raw_word)
            if not words:
                continue
            normalized = " ".join(words)
            normalized_lower = normalized.lower()
            if normalized_lower in {"human agent", "need agents", "agentforce enterprise"}:
                normalized = {
                    "human agent": "Human + agent workflow",
                    "need agents": "AI agents adoption",
                    "agentforce enterprise": "Agentforce enterprise",
                }[normalized_lower]
            root_key = " ".join(word.lower().rstrip("s") for word in words)
            if root_key in seen_roots:
                continue
            seen_roots.add(root_key)
            cleaned_parts.append(normalized)
            if len(cleaned_parts) >= max_terms:
                break

        if not cleaned_parts:
            return "Underserved content theme"
        return " • ".join(word.title() for word in cleaned_parts)

    @staticmethod
    def _company_short_label(name: str, max_len: int = 12) -> str:
        compact_map = {
            "mypromovideos": "MyPromo",
            "salesforce": "Salesforce",
            "mailchimp": "Mailchimp",
            "monday.com": "monday",
            "wistia": "Wistia",
        }
        lowered = (name or "").strip().lower()
        if lowered in compact_map:
            return compact_map[lowered]
        return (name or "")[:max_len]

    @staticmethod
    def _length_label(best: str) -> str:
        best = (best or "").lower()
        if best == "short":
            return "Short"
        if best == "medium":
            return "Medium"
        if best == "long":
            return "Long"
        return "Mixed"

    @staticmethod
    def _length_recommendation(company_analytics: dict) -> str:
        """Return a brief video length recommendation based on analytics."""
        best = company_analytics.get("best_performing_length", "")
        short_count = company_analytics.get("short", {}).get("count", 0)
        medium_count = company_analytics.get("medium", {}).get("count", 0)
        long_count = company_analytics.get("long", {}).get("count", 0)
        total = max(short_count + medium_count + long_count, 1)
        short_share = short_count / total

        if best == "short" and short_share >= 0.8:
            return "Stretch"
        if best == "short":
            return "Expand"
        if best == "medium":
            return "Center"
        if best == "long":
            return "Deepen"
        return "Explore"

    @staticmethod
    def _format_engagement_rate(value: float) -> str:
        if value < 0.1:
            return f"{value:.3f}%"
        return f"{value:.2f}%"

    @staticmethod
    def _duration_label(duration_seconds: int) -> str:
        if duration_seconds <= 0:
            return ""
        minutes, seconds = divmod(duration_seconds, 60)
        return f"{minutes}:{seconds:02d}"

    @staticmethod
    def _published_label(published_at: str) -> str:
        if not published_at:
            return ""
        try:
            dt = datetime.fromisoformat(published_at.replace("Z", "+00:00"))
            return dt.strftime("%d %b %Y")
        except (ValueError, AttributeError):
            return ""

    @staticmethod
    def _video_views_per_day(video: dict) -> str:
        raw_date = video.get("published_at", "")
        try:
            published_at = datetime.fromisoformat(str(raw_date).replace("Z", "+00:00"))
        except (ValueError, AttributeError):
            return ""
        age_days = max((datetime.now(published_at.tzinfo) - published_at).total_seconds() / 86_400, 1.0)
        views = max(int(video.get("view_count", 0)), 0)
        return f"{views / age_days:.1f} views/day"

    @staticmethod
    def _video_url(video: dict) -> str:
        video_id = str(video.get("video_id", "")).strip()
        if not video_id:
            return ""
        return f"https://www.youtube.com/watch?v={video_id}"

    def _gap_excerpt(self, gap_analysis: str, opp_scores: list[dict]) -> str:
        cleaned = self._clean_markdown(gap_analysis or "")
        if cleaned and "unavailable due to Gemini" not in cleaned:
            return self._slide_trim(cleaned, 280)
        if not opp_scores:
            return "Even without one standout gap, the next opportunity is likely to come from stronger proof content, better discovery packaging, or a more deliberate buyer-stage mix."
        top = opp_scores[:3]
        parts = []
        for opp in top:
            parts.append(
                f"{self._humanize_topic_label(opp.get('topic', ''), 2)} scores {opp.get('opportunity_score', 0):.0f}/100 "
                f"with trend interest at {opp.get('trend_interest', 0):.0f} and scarcity at {opp.get('scarcity_score', 0):.0f}."
            )
        return " ".join(parts)

    def _video_performance_hypothesis(self, video: dict) -> str:
        title = str(video.get("title", "") or "").lower()
        views = int(video.get("view_count", 0) or 0)
        likes = int(video.get("like_count", 0) or 0)
        comments = int(video.get("comment_count", 0) or 0)
        duration = int(video.get("duration_seconds", 0) or 0)
        engagement_rate = (likes + comments) / max(views, 1) * 100
        views_per_day_text = self._video_views_per_day(video)
        views_per_day_value = 0.0
        if views_per_day_text:
            try:
                views_per_day_value = float(views_per_day_text.split()[0])
            except (ValueError, IndexError):
                views_per_day_value = 0.0

        if any(token in title for token in ("how to", "guide", "tutorial", "walkthrough")):
            return "Clear search intent and practical education."
        if any(token in title for token in ("case study", "customer", "testimonial", "roi")):
            return "Concrete proof content likely matched high-intent viewers."
        if views_per_day_value >= 100:
            return "Fresh demand appears strong relative to the video's age."
        if duration and duration <= 120:
            return "Fast, easy-to-consume format likely helped early traction."
        if engagement_rate >= 3:
            return "The topic appears to trigger a stronger response than average."
        return "The topic likely matched an active audience need in the sample."

    def _video_actionable_takeaway(self, video: dict) -> str:
        title = str(video.get("title", "") or "").lower()
        duration = int(video.get("duration_seconds", 0) or 0)
        if any(token in title for token in ("how to", "guide", "tutorial", "walkthrough")):
            return "Test another search-led educational video in the same problem area."
        if any(token in title for token in ("case study", "customer", "testimonial", "roi")):
            return "Package a second proof-led story for a similar buyer concern."
        if duration and duration <= 120:
            return "Reuse the short format for another tightly scoped explainer."
        return "Reuse the winning angle with a clearer CTA or follow-up topic."

    def _top_video_slide_summary(self, companies: list[str], company_data: dict) -> str:
        top_videos = []
        for company in companies:
            top_videos.extend(company_data.get(company, {}).get("top_videos", [])[:2])
        if not top_videos:
            return "Use this slide to spot which topics are earning the fastest response in the sample."

        search_intent_hits = 0
        proof_hits = 0
        short_hits = 0
        for video in top_videos:
            title = str(video.get("title", "") or "").lower()
            if any(token in title for token in ("how to", "guide", "tutorial", "walkthrough", "explainer")):
                search_intent_hits += 1
            if any(token in title for token in ("case study", "customer", "testimonial", "roi")):
                proof_hits += 1
            if int(video.get("duration_seconds", 0) or 0) <= 180:
                short_hits += 1

        if search_intent_hits >= max(len(top_videos) // 2, 2):
            return "High-performing videos skew toward educational, search-intent-led topics. The client should test more problem-solving titles with similarly clear intent."
        if proof_hits >= 2:
            return "Winning videos show that proof-led or customer-facing content is earning attention. The client should turn one buyer proof theme into a repeatable series."
        if short_hits >= max(len(top_videos) // 2, 2):
            return "Shorter videos appear to be earning faster traction in this sample. The client should test concise explainers before expanding into longer versions."
        return "Winning videos suggest that clarity of topic matters more than volume alone. The client should double down on the angle with the strongest early response."

    def _channel_overview_takeaway(self, companies: list[str], company_data: dict, target_company: str) -> str:
        target = company_data.get(target_company, {})
        if not target:
            return "Use this slide to understand how much audience scale the client is competing against and where execution efficiency will matter more than sheer size."
        largest = max(companies, key=lambda c: company_data.get(c, {}).get("subscriber_count", 0))
        target_subs = target.get("subscriber_count", 0)
        largest_subs = company_data.get(largest, {}).get("subscriber_count", 0)
        target_er = target.get("avg_engagement_rate", 0.0)
        peer_er_values = [company_data.get(c, {}).get("avg_engagement_rate", 0.0) for c in companies if c != target_company]
        peer_avg_er = sum(peer_er_values) / max(len(peer_er_values), 1)
        if largest != target_company and target_er >= peer_avg_er:
            return (
                f"{target_company} operates with a smaller audience base than {largest}, which has {largest_subs:,} subscribers, "
                f"but {target_company} still delivers {target_er:.3f}% engagement on average. That suggests execution quality may matter more than sheer scale."
            )
        return (
            f"{target_company} currently has {target_subs:,} subscribers versus {largest_subs:,} for the largest peer in this set. "
            "This gap sets the reach challenge, but it also clarifies how much smarter packaging and sequencing need to work for the client."
        )

    def _benchmark_peer_explanation(self, companies: list[str], company_data: dict, rpi_scores: dict, target_company: str) -> tuple[str, str]:
        if target_company not in company_data:
            leader = max(companies, key=lambda company: rpi_scores.get(company, {}).get("rpi_score", 0), default="—")
            return leader, f"{leader} currently sets the strongest overall benchmark in this peer set."

        target = company_data[target_company]
        target_subs = max(float(target.get("subscriber_count", 0)), 1.0)
        target_consistency = float(target.get("consistency_score", 0.0))
        target_er = float(target.get("avg_engagement_rate", 0.0))
        target_funnel = str(target.get("funnel_label", ""))

        best_peer = ""
        best_score = float("inf")
        for company in companies:
            if company == target_company:
                continue
            peer = company_data.get(company, {})
            scale_gap = abs(math.log10(max(float(peer.get("subscriber_count", 0)), 1.0)) - math.log10(target_subs))
            er_gap = abs(float(peer.get("avg_engagement_rate", 0.0)) - target_er)
            consistency_gap = abs(float(peer.get("consistency_score", 0.0)) - target_consistency) / 20.0
            funnel_gap = 0.0 if str(peer.get("funnel_label", "")) == target_funnel else 1.0
            composite_gap = scale_gap + er_gap + consistency_gap + funnel_gap
            if composite_gap < best_score:
                best_score = composite_gap
                best_peer = company

        if not best_peer:
            best_peer = max(companies, key=lambda company: rpi_scores.get(company, {}).get("rpi_score", 0), default="—")
        peer = company_data.get(best_peer, {})
        reason = (
            f"{best_peer} is the most useful benchmark for {target_company} because it combines a comparable audience shape, "
            f"{peer.get('funnel_label', 'mixed')} journey coverage, and {peer.get('avg_engagement_rate', 0.0):.3f}% engagement while outperforming on consistency."
        )
        return best_peer, reason

    def _funnel_takeaway(self, companies: list[str], company_data: dict) -> str:
        if not companies:
            return "Use this slide to see whether recent content mostly serves awareness, evaluation, or proof-stage needs."
        target_company = companies[0]
        target = company_data.get(target_company, {})
        label = target.get("funnel_label", "Balanced")
        return (
            f"{target_company} currently shows {target.get('tofu_pct', 0):.1f}% awareness, {target.get('mofu_pct', 0):.1f}% consideration, "
            f"and {target.get('bofu_pct', 0):.1f}% proof-focused coverage. That {label.lower()} mix suggests how well the channel may support both discovery and decision-stage confidence, and it points to the next buyer-stage experiment."
        )

    def _length_strategy_takeaway(self, companies: list[str], company_data: dict) -> str:
        if not companies:
            return "Use this slide to decide which format length deserves the next test."
        target_company = companies[0]
        target = company_data.get(target_company, {})
        best = target.get("best_performing_length", "mixed")
        confidence = target.get("length_confidence", "LOW")
        return (
            f"{target_company}'s strongest recent format signal appears to be {best} content with {confidence.lower()} confidence. "
            "Use that as the next test direction, but keep in mind that missing or lightly used length buckets should still be tested before making a hard format call."
        )

    def _seo_takeaway(self, companies: list[str], seo_scores: dict) -> str:
        if not companies:
            return "Use this slide to see what is helping or hurting discoverability before viewers even watch."
        leader = max(companies, key=lambda company: seo_scores.get(company, {}).get("seo_score", 0.0))
        weakest = min(companies, key=lambda company: seo_scores.get(company, {}).get("seo_score", 0.0))
        leader_breakdown = seo_scores.get(leader, {}).get("breakdown", {})
        weakest_breakdown = seo_scores.get(weakest, {}).get("breakdown", {})
        weakest_signal = "timestamps" if weakest_breakdown.get("has_timestamps_pct", 0.0) == 0 else "title clarity"
        return (
            f"{leader} has the strongest discovery setup at {seo_scores.get(leader, {}).get('seo_score', 0.0):.1f}/100, led by description depth of {leader_breakdown.get('description_depth', 0.0):.0f}. "
            f"The clearest weakness in this peer set is {weakest}'s {weakest_signal}, so the first SEO fix should focus on making educational content easier to find and navigate."
        )

    def _scorecard_closing(self, companies: list[str], company_data: dict, rpi_scores: dict, seo_scores: dict, target_company: str) -> str:
        target = company_data.get(target_company, {})
        if not target:
            leader = max(companies, key=lambda company: rpi_scores.get(company, {}).get("rpi_score", 0.0), default="—")
            return f"{leader} currently leads the peer set, but the point of this scorecard is to surface the specific habits the client can borrow next."
        strength = "engagement efficiency" if target.get("avg_engagement_rate", 0.0) >= 1.0 else "audience focus"
        biggest_opportunity = "proof-stage content" if target.get("bofu_pct", 0.0) < 15 else "publishing consistency"
        return (
            f"{target_company} already shows a credible strength in {strength}. The biggest opportunity is {biggest_opportunity}, "
            "which means the client does not need a full reset, only a sharper next set of experiments."
        )

    def _action_plan_takeaway(self, target_company: str, target_data: dict, recommendations: list[dict], opp_scores: list[dict]) -> str:
        whitespace = self._humanize_topic_label(opp_scores[0].get("topic", ""), 2) if opp_scores else "the strongest whitespace theme"
        top_move = recommendations[0].get("title", "the top-priority recommendation") if recommendations else "the top-priority recommendation"
        return (
            f"This plan turns {top_move.lower()} into a practical 90-day sequence for {target_company}. "
            f"It is grounded in the current {target_data.get('mean_gap_days', 0):.1f}-day cadence, the {target_data.get('funnel_label', 'current')} journey mix, and whitespace around {whitespace}."
        )

    def _action_plan_sections(self, text: str) -> dict[str, str]:
        cleaned = self._clean_markdown(text)
        patterns = {
            "Week 1–2": r"Week 1[–-]2:\s*(.*?)(?=Month 1:|Month 2[–-]3:|$)",
            "Month 1": r"Month 1:\s*(.*?)(?=Month 2[–-]3:|$)",
            "Month 2–3": r"Month 2[–-]3:\s*(.*)$",
        }
        sections: dict[str, str] = {}
        for label, pattern in patterns.items():
            match = re.search(pattern, cleaned, flags=re.IGNORECASE)
            if match:
                sections[label] = " ".join(match.group(1).split())
        return sections

    def _action_plan_card_bodies(
        self,
        phases: dict[str, str],
        recommendations: list[dict],
        target_company: str,
        target_data: dict,
        opp_scores: list[dict],
    ) -> list[str]:
        whitespace = self._humanize_topic_label(opp_scores[0].get("topic", ""), 2) if opp_scores else "highest-value theme"
        top_titles = [self._slide_trim(rec.get("title", ""), 30) for rec in recommendations[:2] if rec.get("title")]
        first_focus = top_titles[0] if top_titles else "top-priority move"
        second_focus = top_titles[1] if len(top_titles) > 1 else "packaging upgrade"
        cadence = self._format_gap_days(float(target_data.get("mean_gap_days", 0.0) or 0.0))
        best_length = self._length_label(target_data.get("best_performing_length", "mixed"))

        defaults = {
            "Week 1–2": (
                f"Focus: {first_focus}\n"
                f"Do: confirm one theme around {whitespace}, tighten title and description structure, and brief the next 3 videos.\n"
                f"Measure: baseline views/day and current {cadence} publishing gap."
            ),
            "Month 1": (
                f"Focus: launch the first test tied to {second_focus.lower()}.\n"
                "Do: publish one proof-led or educational asset, then review early engagement and click quality after 14 days.\n"
                "Measure: views/day, engagement rate, and completion quality."
            ),
            "Month 2–3": (
                f"Focus: turn the best signal into a repeatable {best_length.lower()} series.\n"
                "Do: keep the strongest angle, remove the weakest experiment, and tighten cadence around a predictable schedule.\n"
                "Measure: repeat-viewer response and consistency score trend."
            ),
        }

        ordered = []
        for label in ["Week 1–2", "Month 1", "Month 2–3"]:
            source = self._clean_markdown(phases.get(label, ""))
            if source:
                summary = self._slide_trim(source, 150)
                if label == "Week 1–2":
                    ordered.append(
                        f"Focus: {first_focus}\n"
                        f"Do: {summary}\n"
                        f"Measure: baseline views/day and current {cadence} publishing gap."
                    )
                elif label == "Month 1":
                    ordered.append(
                        f"Focus: {second_focus}\n"
                        f"Do: {summary}\n"
                        "Measure: 14-day engagement and early views/day."
                    )
                else:
                    ordered.append(
                        f"Focus: scale what worked\n"
                        f"Do: {summary}\n"
                        "Measure: repeat-view response and publishing discipline."
                    )
            else:
                ordered.append(defaults[label])
        return ordered

    @staticmethod
    def _rpi_grade(rpi_score: float) -> str:
        """Convert RPI score to an A/B/C/D grade."""
        if rpi_score >= 75:
            return "A"
        if rpi_score >= 50:
            return "B"
        if rpi_score >= 25:
            return "C"
        return "D"
